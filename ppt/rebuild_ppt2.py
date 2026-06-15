# -*- coding: utf-8 -*-
"""Rebuild 02_演算法設計.pptx  – 5-section structure, Algorithm focus."""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.chart.data import ChartData
from pptx.enum.chart import XL_CHART_TYPE

BG    = RGBColor(0x0F,0x0A,0x1E)
PANEL = RGBColor(0x1A,0x12,0x30)
PURP  = RGBColor(0x9B,0x5D,0xE5)
PURPL = RGBColor(0xC7,0x7D,0xFF)
TEAL  = RGBColor(0x00,0xF5,0xD4)
GOLD  = RGBColor(0xFE,0xE4,0x40)
WHITE = RGBColor(0xF0,0xEA,0xF8)
DIM   = RGBColor(0x8B,0x80,0xA8)
RED   = RGBColor(0xF1,0x5B,0xB5)
BLUE  = RGBColor(0x00,0xBB,0xF9)
GREEN = RGBColor(0x57,0xCC,0x99)
CODE  = RGBColor(0x0A,0x07,0x18)
DARK2 = RGBColor(0x13,0x0D,0x22)

prs = Presentation()
prs.slide_width  = Inches(10)
prs.slide_height = Inches(5.625)

def bg(sl): sl.background.fill.solid(); sl.background.fill.fore_color.rgb = BG

def box(sl,text,x,y,w,h,sz=14,bold=False,color=WHITE,
        align=PP_ALIGN.LEFT,italic=False,font="Calibri"):
    tb=sl.shapes.add_textbox(Inches(x),Inches(y),Inches(w),Inches(h))
    tf=tb.text_frame; tf.word_wrap=True
    p=tf.paragraphs[0]; p.alignment=align
    r=p.add_run(); r.text=text
    r.font.size=Pt(sz); r.font.bold=bold; r.font.italic=italic
    r.font.color.rgb=color; r.font.name=font

def rect(sl,x,y,w,h,fill=PANEL,line=None,radius=False):
    st=5 if radius else 1
    sh=sl.shapes.add_shape(st,Inches(x),Inches(y),Inches(w),Inches(h))
    sh.fill.solid(); sh.fill.fore_color.rgb=fill
    if line: sh.line.color.rgb=line; sh.line.width=Pt(0.5)
    else: sh.line.fill.background()
    return sh

def oval(sl,x,y,w,h,fill):
    sh=sl.shapes.add_shape(9,Inches(x),Inches(y),Inches(w),Inches(h))
    sh.fill.solid(); sh.fill.fore_color.rgb=fill; sh.line.fill.background()

def tag(sl,text):
    rect(sl,0.4,0.15,2.8,0.38,fill=PURP)
    box(sl,text,0.4,0.15,2.8,0.38,sz=11,bold=True,color=WHITE,align=PP_ALIGN.CENTER)

def title(sl,text):
    box(sl,text,0.4,0.62,9.2,0.72,sz=28,bold=True,color=WHITE)

def stat_card(sl,val,lbl,x,y,vc=TEAL):
    rect(sl,x,y,2.05,0.88)
    box(sl,val,x,y+0.04,2.05,0.46,sz=22,bold=True,color=vc,align=PP_ALIGN.CENTER)
    box(sl,lbl,x,y+0.52,2.05,0.3,sz=10,color=DIM,align=PP_ALIGN.CENTER)

# ════════════════════════════════════════════════════════
# S1 – Title
# ════════════════════════════════════════════════════════
sl=prs.slides.add_slide(prs.slide_layouts[6]); bg(sl)
for i in range(10):
    r=1.2+i*0.5
    sh=sl.shapes.add_shape(9,Inches(5-r/2),Inches(2.8-r/2),Inches(r),Inches(r))
    sh.fill.background()
    sh.line.color.rgb=RGBColor(80+i*10,20+i*5,150+i*8); sh.line.width=Pt(1)

rect(sl,0.4,0.38,3.2,0.5,fill=PURP)
box(sl,"BULLET STORM",0.4,0.38,3.2,0.5,sz=14,bold=True,color=WHITE,align=PP_ALIGN.CENTER)
box(sl,"演算法設計",0.4,1.1,9,1.1,sz=48,bold=True,color=WHITE)
box(sl,"Algorithm Design & Internal Mechanics",0.4,2.35,8,0.5,sz=17,color=PURPL)
box(sl,"Project Motivation  ·  System Architecture  ·  Technologies Used  ·  Experimental Results  ·  Future Improvements",
    0.4,3.05,9.2,0.4,sz=11,color=DIM)
for i,(v,l) in enumerate([("8","彈幕演算法"),("4","運動類型"),("5","技能樹"),("Lv 1-99","EXP 系統")]):
    stat_card(sl,v,l,0.4+i*2.35,3.65)

# ════════════════════════════════════════════════════════
# S2 – Project Motivation
# ════════════════════════════════════════════════════════
sl=prs.slides.add_slide(prs.slide_layouts[6]); bg(sl)
tag(sl,"01  Project Motivation")
title(sl,"專案動機")

motivations=[
    ("∿","數學密集型運算挑戰",
     "三角函數、向量、參數方程式在每幀執行數百次——60fps × 400 顆子彈 = 每秒 24,000 次物理計算"),
    ("⊕","碰撞偵測的效能壓力",
     "子彈 × 玩家 × 敵人的 O(n²) 判定，需透過早期終止與延遲清除讓每幀維持在 1ms 以內"),
    ("AI","敵人行為的演算法設計",
     "入場路徑、巡邏 AI、多階段射擊時機，全都需要有狀態的決策邏輯驅動"),
    ("||","難度的多維度平衡",
     "速度、數量、頻率三個軸同步調整，設計出難度曲線平滑遞增、不讓玩家感覺突然跳難的系統"),
]
for i,(ic,h,d) in enumerate(motivations):
    y=1.52+i*0.96
    rect(sl,0.4,y,9.2,0.82)
    oval(sl,0.55,y+0.25,0.32,0.32,PURP)
    box(sl,ic,0.55,y+0.25,0.32,0.32,sz=14,align=PP_ALIGN.CENTER)
    box(sl,h,1.0,y+0.08,8.5,0.28,sz=13,bold=True,color=WHITE)
    box(sl,d,1.0,y+0.38,8.5,0.32,sz=11,color=DIM)

box(sl,"目標：僅用標準 Java，以可組合的靜態演算法庫實現 8 種彈幕 × 4 種物理 × 5 種技能升級",
    0.4,5.2,9.2,0.28,sz=10,color=DIM,italic=True)

# ════════════════════════════════════════════════════════
# S3 – System Architecture
# ════════════════════════════════════════════════════════
sl=prs.slides.add_slide(prs.slide_layouts[6]); bg(sl)
tag(sl,"02  System Architecture")
title(sl,"核心類別架構與物件關係")

classes=[
    ("GamePanel","主迴圈 60fps\n遊戲狀態管理",PURP,  3.9, 1.52),
    ("BulletPattern","8 種演算法庫\n(static，無狀態)",TEAL, 0.4, 2.75),
    ("Boss","5 階段 Boss AI\n呼叫 BulletPattern",RED,  3.55,2.75),
    ("Enemy","3 類敵人 + AI\n自動生成與消失",BLUE, 6.8, 2.75),
    ("Bullet","通用子彈 type 0-3\n所有子彈共用",GOLD, 0.4, 3.95),
    ("Player","玩家控制\n整合 EXP 系統",GREEN,3.55,3.95),
    ("ExperienceSystem","5 項技能\n自動/手動升級",PURPL,6.8, 3.95),
]
for (nm,role,c,x,y) in classes:
    rect(sl,x,y,2.5,0.85)
    rect(sl,x+0.08,y+0.07,2.34,0.32,fill=c)
    box(sl,nm,x+0.08,y+0.07,2.34,0.32,sz=11,bold=True,color=BG,align=PP_ALIGN.CENTER)
    box(sl,role,x+0.1,y+0.44,2.3,0.36,sz=9.5,color=WHITE,align=PP_ALIGN.CENTER)

# Vertical lines
for ax in [1.65,4.8,8.05]:
    sh=sl.shapes.add_shape(1,Inches(ax),Inches(3.6),Inches(0.02),Inches(0.35))
    sh.fill.solid(); sh.fill.fore_color.rgb=PURP; sh.line.fill.background()
# Horizontal lines
for ax in [3.3,6.45]:
    sh=sl.shapes.add_shape(1,Inches(ax),Inches(2.375),Inches(0.25),Inches(0.02))
    sh.fill.solid(); sh.fill.fore_color.rgb=PURP; sh.line.fill.background()

box(sl,"BulletPattern 為純靜態工具類（無實例化），Enemy 與 Boss 傳入參數即回傳 List<Bullet>，達成演算法複用與解耦。",
    0.4,5.05,9.2,0.35,sz=10,color=DIM,italic=True)

# ════════════════════════════════════════════════════════
# S4 – Technologies: 8 Bullet Algorithms
# ════════════════════════════════════════════════════════
sl=prs.slides.add_slide(prs.slide_layouts[6]); bg(sl)
tag(sl,"03  Technologies Used")
title(sl,"8 種彈幕生成演算法")

patterns=[
    ("circularBurst()","環形爆散","全向 n 發","θᵢ = θ₀ + i×(2π/n)",PURPL),
    ("aimedSpread()",  "瞄準扇射","鎖定玩家扇形","base = atan2(Δy,Δx) ± spread",TEAL),
    ("doubleRing()",   "雙環交轉","外快內慢對向","outer_v > 0, inner_v < 0",GOLD),
    ("spiralArm()",    "螺旋臂","每幀遞增旋轉","offset += Δθ per frame",BLUE),
    ("homingBullets()","追蹤彈","自動追蹤玩家","angle = lerp(angle, target, 0.05)",RED),
    ("waveBullets()",  "波動彈","蛇形軌跡","offset = sin(phase) × amp",GREEN),
    ("starBurst()",    "星爆交錯","交替速度環","i%2==0→fast, i%2==1→slow",PURPL),
    ("butterflySpray()","蝴蝶散射","蝴蝶型 8 字","spread = sin(t) × π × 0.4",TEAL),
]
for i,(fn,zh,desc,formula,c) in enumerate(patterns):
    col,row=i%2,i//2
    x=0.4+col*4.85; y=1.52+row*0.98
    rect(sl,x,y,4.65,0.84)
    oval(sl,x+0.1,y+0.08,0.28,0.28,c)
    box(sl,zh,x+0.5,y+0.04,1.4,0.28,sz=12,bold=True,color=WHITE)
    box(sl,fn, x+0.5,y+0.34,1.6,0.22,sz=8,color=DIM,font="Courier New")
    box(sl,formula,x+2.2,y+0.04,2.35,0.28,sz=9,color=GOLD,font="Courier New")
    box(sl,desc,  x+2.2,y+0.34,2.35,0.22,sz=9,color=DIM)

# ════════════════════════════════════════════════════════
# S5 – Technologies: Bullet Physics + EXP
# ════════════════════════════════════════════════════════
sl=prs.slides.add_slide(prs.slide_layouts[6]); bg(sl)
tag(sl,"03  Technologies Used")
title(sl,"子彈物理模式與 EXP 升級系統")

# Left: 4 movement types
box(sl,"4 種子彈運動模式（Bullet.type 欄位）",0.4,1.52,4.55,0.32,sz=13,bold=True,color=PURPL)
types=[
    ("type 0","直線",BLUE,  ["x += vx","y += vy"],"等速直線"),
    ("type 1","追蹤",RED,   ["angle → atan2(target)","clamp(Δangle, ±ω)"],"漸進轉向"),
    ("type 2","波動",TEAL,  ["phase += freq","perp = sin(phase)×amp"],"正弦側移"),
    ("type 3","加速",GOLD,  ["speed += accel","vx,vy 依 speed 更新"],"可負 accel"),
]
for i,(tp,nm,c,code,note) in enumerate(types):
    y=1.92+i*0.9
    rect(sl,0.4,y,4.55,0.76)
    rect(sl,0.5,y+0.08,0.7,0.3,fill=c)
    box(sl,tp,0.5,y+0.08,0.7,0.3,sz=9.5,bold=True,color=BG,align=PP_ALIGN.CENTER)
    box(sl,nm,1.3,y+0.06,0.8,0.3,sz=12,color=WHITE)
    rect(sl,2.2,y+0.06,2.6,0.64,fill=CODE)
    for j,ln in enumerate(code):
        box(sl,ln,2.28,y+0.1+j*0.26,2.45,0.22,sz=8.5,color=TEAL,font="Courier New")
    box(sl,note,2.2,y+0.5,2.6,0.22,sz=8,color=DIM)

# Right: EXP system
box(sl,"EXP 升級系統（5 項技能樹）",5.15,1.52,4.45,0.32,sz=13,bold=True,color=PURPL)
formulas=[
    ("升級所需","expNext = 1000 + level × 500",TEAL),
    ("擊殺 EXP","50 × (type+1) × (diff+1)",GREEN),
    ("Boss EXP", "5000 × (diff+1)",GOLD),
    ("刷掠 EXP", "5 / bullet",BLUE),
    ("射速技能",  "speed × (1.0 + lvl×0.15)",PURPL),
    ("射頻技能",  "cooldown -= lvl×2  (min 2)",RED),
    ("生命技能",  "max_life += lvl",GREEN),
    ("炸彈技能",  "bomb_dmg × (1+lvl×0.2)",GOLD),
    ("刷掠技能",  "graze_r × (1+lvl×0.1)",TEAL),
]
for i,(lbl,f,c) in enumerate(formulas):
    y=1.92+i*0.4
    box(sl,lbl+"：",5.15,y,1.5,0.32,sz=9,color=DIM)
    rect(sl,6.68,y-0.02,2.88,0.34,fill=CODE)
    box(sl,f,6.75,y,2.76,0.3,sz=8.5,color=c,font="Courier New")

# ════════════════════════════════════════════════════════
# S6 – Technologies: Enemy AI
# ════════════════════════════════════════════════════════
sl=prs.slides.add_slide(prs.slide_layouts[6]); bg(sl)
tag(sl,"03  Technologies Used")
title(sl,"敵人 AI 生命週期管理")

# Lifecycle pipeline
stages=[
    ("生成","age = -(i×15)\n錯開入場",GREEN),
    ("進場","age < 60\n依 vx/vy 飛入",BLUE),
    ("減速","60 ≤ age < 120\nvx *= 0.95",TEAL),
    ("巡邏","age ≥ 120\nsin(t)×40 擺動",GOLD),
    ("消失","age ≥ maxAge\nactive = false",RED),
]
for i,(nm,desc,c) in enumerate(stages):
    x=0.4+i*1.16
    rect(sl,x,1.52,1.06,1.05)
    oval(sl,x+0.31,1.62,0.44,0.44,c)
    box(sl,nm,x+0.31,1.62,0.44,0.44,sz=9,bold=True,color=BG,align=PP_ALIGN.CENTER)
    box(sl,desc,x+0.04,2.18,0.98,0.36,sz=8.5,color=DIM,align=PP_ALIGN.CENTER,font="Courier New")
    if i<4:
        sh=sl.shapes.add_shape(1,Inches(x+1.06),Inches(2.0),Inches(0.10),Inches(0.02))
        sh.fill.solid();sh.fill.fore_color.rgb=PURP;sh.line.fill.background()

# maxAge table
box(sl,"各敵人存活時間",0.4,2.75,5.7,0.3,sz=13,bold=True,color=WHITE)
rows=[
    ["敵人類型","maxAge (幀)","存活秒數","射擊間隔"],
    ["type 0 仙子","480","8.0 秒","100 幀"],
    ["type 1 妖怪","600","10.0 秒","90 幀"],
    ["type 2 頭目","720","12.0 秒","70 幀"],
]
tbl=sl.shapes.add_table(4,4,Inches(0.4),Inches(3.1),Inches(5.7),Inches(1.5)).table
cw=[1.5,1.2,1.5,1.0]
for ci,w in enumerate(cw): tbl.columns[ci].width=Inches(w)
for ri,row in enumerate(rows):
    for ci,txt in enumerate(row):
        cell=tbl.cell(ri,ci); cell.text=txt
        p=cell.text_frame.paragraphs[0]; p.alignment=PP_ALIGN.CENTER
        run=p.runs[0]; run.font.size=Pt(11); run.font.name="Calibri"
        run.font.bold=(ri==0); run.font.color.rgb=BG if ri==0 else WHITE
        cell.fill.solid(); cell.fill.fore_color.rgb=PURP if ri==0 else PANEL

# Code block
rect(sl,6.3,1.52,3.3,3.08,fill=DARK2)
box(sl,"錯開生成演算法",6.4,1.6,3.1,0.32,sz=12,bold=True,color=TEAL,align=PP_ALIGN.CENTER)
code_lines=[
    ("for (int i = 0; i < n; i++) {",False),
    ("  Enemy en = new Enemy(...);",False),
    ("  en.age = -(i * 15);",False),
    ("  // 每隻差 15 幀入場",True),
    ("  enemies.add(en);",False),
    ("}",False),
    ("// update() 中",True),
    ("if (age <= 0) return;",False),
    ("if (age >= maxAge) {",False),
    ("  active = false;",False),
    ("}",False),
]
for j,(ln,cmt) in enumerate(code_lines):
    box(sl,ln,6.4,2.0+j*0.21,3.12,0.19,sz=8.2,
        color=DIM if cmt else TEAL,font="Courier New")

# ════════════════════════════════════════════════════════
# S7 – Experimental Results
# ════════════════════════════════════════════════════════
sl=prs.slides.add_slide(prs.slide_layouts[6]); bg(sl)
tag(sl,"04  Experimental Results")
title(sl,"實驗結果：效能與難度平衡")

# Bullet count per difficulty chart
box(sl,"各難度每秒產生子彈數量",0.4,1.52,5.5,0.32,sz=13,bold=True,color=WHITE)
cd=ChartData()
cd.categories=["Easy","Normal","Hard","Lunatic"]
cd.add_series("一般敵人",[1.2,1.7,2.2,2.8])
cd.add_series("Boss 戰", [2.5,3.8,5.2,7.0])
ch=sl.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED,
    Inches(0.4),Inches(1.9),Inches(5.5),Inches(2.55),cd)
ch.chart.plots[0].series[0].format.fill.solid()
ch.chart.plots[0].series[0].format.fill.fore_color.rgb=BLUE
ch.chart.plots[0].series[1].format.fill.solid()
ch.chart.plots[0].series[1].format.fill.fore_color.rgb=RED

# Metrics right
metrics=[
    ("~ 400 顆","最大同時子彈數",GOLD),
    ("< 1 ms","碰撞偵測每幀耗時",GREEN),
    ("60 FPS","幀率穩定度 ±1%",TEAL),
    ("< 80 MB","記憶體占用",BLUE),
]
for i,(v,lbl,c) in enumerate(metrics):
    y=1.52+i*0.98
    rect(sl,6.1,y,3.5,0.84)
    sh=sl.shapes.add_shape(1,Inches(6.1),Inches(y),Inches(0.07),Inches(0.84))
    sh.fill.solid();sh.fill.fore_color.rgb=c;sh.line.fill.background()
    box(sl,lbl,6.28,y+0.04,3.2,0.26,sz=10,color=DIM)
    box(sl,v,  6.28,y+0.3, 3.2,0.42,sz=18,bold=True,color=c)

box(sl,"難度公式：speedMult = (1.0 + difficulty×0.2)×0.65  |  Easy 0.65× / Normal 0.78× / Hard 0.91× / Lunatic 1.04×",
    0.4,4.58,9.2,0.38,sz=8.5,color=DIM,italic=True)

# ════════════════════════════════════════════════════════
# S8 – Future Improvements
# ════════════════════════════════════════════════════════
sl=prs.slides.add_slide(prs.slide_layouts[6]); bg(sl)
tag(sl,"05  Future Improvements")
title(sl,"未來改進方向")

phases=[
    ("短期","Short-term",GREEN,[
        ("空間網格分割 (Grid)","碰撞偵測由 O(n²) 降至 O(n log n)"),
        ("物件池 (Object Pool)","減少 Bullet 物件 GC 壓力"),
        ("Bezier 曲線彈道","更自然的參數化子彈路徑"),
    ]),
    ("中期","Mid-term",TEAL,[
        ("難度自適應","依存活率動態調整子彈密度"),
        ("Boss AI 狀態機","Idle / Attack / Enrage 三態切換"),
        ("10+ 種技能樹","更深度的 RPG 式成長系統"),
    ]),
    ("長期","Long-term",PURPL,[
        ("強化學習難度 (RL)","以 RL 訓練最佳難度曲線"),
        ("程序生成 Boss 技能","隨機但平衡的攻擊組合"),
        ("多人彈幕同步","分散式狀態一致性演算法"),
    ]),
]
for i,(zh,en,c,items) in enumerate(phases):
    x=0.4+i*3.22
    rect(sl,x,1.52,3.05,3.75)
    rect(sl,x+0.1,1.62,2.85,0.48,fill=c)
    box(sl,f"{zh}  {en}",x+0.1,1.62,2.85,0.48,sz=13,bold=True,color=BG,align=PP_ALIGN.CENTER)
    for j,(h,d) in enumerate(items):
        y=2.22+j*1.0
        oval(sl,x+0.2,y+0.05,0.22,0.22,c)
        box(sl,h,x+0.52,y+0.02,2.4,0.28,sz=11,bold=True,color=WHITE)
        box(sl,d,x+0.52,y+0.3,2.4,0.3,sz=10,color=DIM)

# ════════════════════════════════════════════════════════
# S9 – Closing
# ════════════════════════════════════════════════════════
sl=prs.slides.add_slide(prs.slide_layouts[6]); bg(sl)
for i in range(10):
    r=0.8+i*0.55
    sh=sl.shapes.add_shape(9,Inches(5-r/2),Inches(2.8-r/2),Inches(r),Inches(r))
    sh.fill.background()
    sh.line.color.rgb=RGBColor(60+i*10,10+i*5,120+i*10); sh.line.width=Pt(0.9)

box(sl,"BULLET STORM",0.5,0.75,9,0.8,sz=42,bold=True,color=PURPL,align=PP_ALIGN.CENTER)
box(sl,"演算法設計篇",0.5,1.62,9,0.45,sz=18,color=PURP,align=PP_ALIGN.CENTER)
for i,(v,l,s) in enumerate([
    ("8","彈幕演算法","從環形到追蹤"),
    ("4","運動物理","直線到加速"),
    ("5","技能樹","射速到刷掠"),
    ("99","最高等級","線性 EXP 成長"),
]):
    x=0.5+i*2.28
    rect(sl,x,2.85,2.1,1.8)
    box(sl,v,x,2.93,2.1,0.62,sz=26,bold=True,color=TEAL,align=PP_ALIGN.CENTER)
    box(sl,l,x,3.56,2.1,0.28,sz=11,bold=True,color=WHITE,align=PP_ALIGN.CENTER)
    box(sl,s,x,3.86,2.1,0.55,sz=9,color=DIM,align=PP_ALIGN.CENTER)

box(sl,"github.com/Shirodot/BulletStorm",0.5,4.9,9,0.35,sz=11,color=DIM,align=PP_ALIGN.CENTER)

prs.save(r"C:\project T\BulletStorm\ppt\02_演算法設計.pptx")
print("PPT2 saved")
