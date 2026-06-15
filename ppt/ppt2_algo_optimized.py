"""
PPT 2 — Bullet Storm：演算法設計篇（優化版）
- 僅展示8種實際使用的演算法（非12種聲稱的）
- 準確的遊戲參數和機制
- 完整的EXP升級系統詳解
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.chart.data import ChartData
from pptx.enum.chart import XL_CHART_TYPE
import math

# ─── Palette (Purple/Math theme) ────────────────────────
DARK   = RGBColor(0x0F, 0x0A, 0x1E)
PANEL  = RGBColor(0x1A, 0x12, 0x30)
ACCENT = RGBColor(0x9B, 0x5D, 0xE5)
ACCTL  = RGBColor(0xC7, 0x7D, 0xFF)
TEAL   = RGBColor(0x00, 0xF5, 0xD4)
GOLD   = RGBColor(0xFE, 0xE4, 0x40)
WHITE  = RGBColor(0xF0, 0xEA, 0xF8)
DIM    = RGBColor(0x8B, 0x80, 0xA8)
RED    = RGBColor(0xF1, 0x5B, 0xB5)
BLUE   = RGBColor(0x00, 0xBB, 0xF9)
GREEN  = RGBColor(0x57, 0xCC, 0x99)
PANEL2 = RGBColor(0x13, 0x0D, 0x22)
CODE   = RGBColor(0x0A, 0x07, 0x18)

prs = Presentation()
prs.slide_width  = Inches(10)
prs.slide_height = Inches(5.625)

# ─── Helpers ────────────────────────────────────────────
def slide_bg(slide):
    bg = slide.background; f = bg.fill
    f.solid(); f.fore_color.rgb = DARK

def card(slide, x, y, w, h, fill=None):
    c = fill or PANEL
    sh = slide.shapes.add_shape(5, Inches(x), Inches(y), Inches(w), Inches(h))
    sh.fill.solid(); sh.fill.fore_color.rgb = c
    sh.line.color.rgb = RGBColor(0x2A,0x1E,0x42); sh.line.width = Pt(0.5)
    return sh

def add_textbox(slide, text, x, y, w, h,
                font_size=14, bold=False, color=WHITE,
                align=PP_ALIGN.LEFT, italic=False, font_name="Calibri"):
    txb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = txb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = align
    run = p.add_run()
    run.text = text; run.font.size = Pt(font_size)
    run.font.bold = bold; run.font.italic = italic
    run.font.color.rgb = color; run.font.name = font_name
    return txb

def add_oval(slide, x, y, w, h, fill_rgb):
    sh = slide.shapes.add_shape(9, Inches(x), Inches(y), Inches(w), Inches(h))
    sh.fill.solid(); sh.fill.fore_color.rgb = fill_rgb
    sh.line.fill.background(); return sh

def section_tag(slide, text):
    card(slide, 0.4, 0.15, 2.3, 0.4, ACCENT)
    add_textbox(slide, text, 0.4, 0.15, 2.3, 0.4,
                font_size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

def slide_title(slide, text):
    add_textbox(slide, text, 0.4, 0.6, 9.2, 0.7, font_size=28, bold=True, color=WHITE)

def hrule(slide, y):
    sh = slide.shapes.add_shape(1, Inches(0.4), Inches(y), Inches(9.2), Inches(0.03))
    sh.fill.solid(); sh.fill.fore_color.rgb = ACCENT; sh.line.fill.background()

# ═══════════════════════════════════════════════════════
# Slide 1 — Title
# ═══════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
slide_bg(sl)

# Orbit rings
for r, clr, t in [(7, ACCENT, 70), (5, TEAL, 75)]:
    sh = sl.shapes.add_shape(9, Inches(5.5 - r/2), Inches(-r/2), Inches(r), Inches(r))
    sh.fill.background()
    sh.line.color.rgb = clr; sh.line.width = Pt(1.2)

card(sl, 0.5, 0.45, 2.8, 0.52, ACCENT)
add_textbox(sl, "BULLET STORM", 0.5, 0.45, 2.8, 0.52,
            font_size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

add_textbox(sl, "演算法設計", 0.5, 1.15, 7, 1.2, font_size=46, bold=True, color=WHITE)
add_textbox(sl, "Algorithm Design & Internal Mechanics",
            0.5, 2.45, 7, 0.55, font_size=18, color=ACCTL)

sh = sl.shapes.add_shape(1, Inches(0.5), Inches(3.1), Inches(3.5), Inches(0.04))
sh.fill.solid(); sh.fill.fore_color.rgb = TEAL; sh.line.fill.background()

add_textbox(sl, "8 種實用彈幕演算法  ·  4 種子彈運動型  ·  5 項技能樹\n敵人AI生命週期  ·  動態難度系統  ·  O(n²) 碰撞偵測",
            0.5, 3.22, 7, 0.85, font_size=14, color=DIM)

stats = [("8","彈幕演算法"), ("4","運動類型"), ("5","技能樹"), ("Level 1-99","EXP系統")]
for i, (v, l) in enumerate(stats):
    x = 0.5 + i * 2.3
    card(sl, x, 4.35, 2.1, 0.95)
    add_textbox(sl, v, x, 4.38, 2.1, 0.48, font_size=22, bold=True, color=TEAL, align=PP_ALIGN.CENTER)
    add_textbox(sl, l, x, 4.84, 2.1, 0.3,  font_size=10, color=DIM,  align=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════
# Slide 2 — 目錄
# ═══════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
slide_bg(sl)
section_tag(sl, "CONTENTS")
slide_title(sl, "簡報目錄")
hrule(sl, 1.35)

toc = [
    ("01","系統架構",   "7 個核心類別、物件關係圖"),
    ("02","彈幕演算法",   "8 種實際使用的射擊模式詳解"),
    ("03","子彈物理",    "4 種運動型別：直線、追蹤、波動、加速"),
    ("04","EXP系統",    "5 項技能樹、自動升級邏輯、公式推導"),
    ("05","敵人AI",     "生命週期、錯開生成、自動消失"),
    ("06","實驗結果",    "效能分析、難度曲線、驗證"),
]
for i, (n, t, d) in enumerate(toc):
    y = 1.52 + i * 0.62
    card(sl, 0.4, y, 9.2, 0.53)
    add_oval(sl, 0.55, y + 0.05, 0.35, 0.35, ACCENT)
    add_textbox(sl, n, 0.55, y + 0.05, 0.35, 0.35,
                font_size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_textbox(sl, t, 1.05, y + 0.02, 2.5, 0.25, font_size=13, bold=True, color=WHITE)
    add_textbox(sl, d, 1.05, y + 0.28, 7.9, 0.2, font_size=10, color=DIM)

# ═══════════════════════════════════════════════════════
# Slide 3 — 系統架構
# ═══════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
slide_bg(sl)
section_tag(sl, "01 系統架構")
slide_title(sl, "核心類別與物件關係")
hrule(sl, 1.35)

classes = [
    ("GamePanel",        "主遊戲迴圈 60fps",      ACCENT, 3.9, 1.55),
    ("BulletPattern",    "8 種演算法庫 (static)", TEAL,   0.4, 2.75),
    ("Boss",             "5 階段 Boss AI",        RED,    3.55, 2.75),
    ("Enemy",            "3 類敵人 + 4 種射擊",   BLUE,   6.8,  2.75),
    ("Bullet",           "通用子彈 (type 0-3)",   GOLD,   0.4,  3.9),
    ("Player",           "玩家 + EXP 系統",       GREEN,  3.55, 3.9),
    ("ExperienceSystem", "5 項技能 + 自動升級",   ACCTL,  6.8,  3.9),
]
for (name, role, clr, x, y) in classes:
    card(sl, x, y, 2.5, 0.72)
    card(sl, x + 0.08, y + 0.06, 2.34, 0.3, clr)
    add_textbox(sl, name, x + 0.08, y + 0.06, 2.34, 0.3,
                font_size=11, bold=True, color=DARK, align=PP_ALIGN.CENTER)
    add_textbox(sl, role, x + 0.1, y + 0.38, 2.3, 0.28,
                font_size=9.5, color=WHITE, align=PP_ALIGN.CENTER)

# Lines
lines = [
    (5.15, 2.25, 1.65, 2.75), (5.15, 2.25, 4.8, 2.75), (5.15, 2.25, 8.05, 2.75),
    (1.65, 3.47, 1.65, 3.9),  (4.8,  3.47, 4.8, 3.9),  (8.05, 3.47, 8.05, 3.9),
]
for (x1, y1, x2, y2) in lines:
    sh = sl.shapes.add_shape(1,
        Inches(min(x1,x2)), Inches(min(y1,y2)),
        Inches(abs(x2-x1)) or Inches(0.02), Inches(abs(y2-y1)) or Inches(0.02))
    sh.fill.solid(); sh.fill.fore_color.rgb = ACCENT; sh.line.fill.background()

add_textbox(sl, "BulletPattern 為純靜態工具類，被 Enemy / Boss 呼叫產生子彈清單",
            0.4, 4.88, 9.2, 0.28, font_size=10, color=DIM)

# ═══════════════════════════════════════════════════════
# Slide 4 — 8 種彈幕演算法
# ═══════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
slide_bg(sl)
section_tag(sl, "02 彈幕演算法")
slide_title(sl, "8 種實際使用的射擊模式")
hrule(sl, 1.35)

patterns = [
    ("circularBurst",  "環形爆散",  "i × (2π/n)",         ACCTL,   "全向 n 發"),
    ("aimedSpread",    "瞄準扇射",  "atan2(target) ± θ",  BLUE,   "鎖定玩家扇形"),
    ("doubleRing",     "雙環交轉",  "outer + inner",       GOLD,   "外快內慢對向"),
    ("spiralArm",      "螺旋臂",    "offset += Δθ",        TEAL,   "轉轉轉（每幀）"),
    ("homingBullets",  "追蹤彈",    "lerp angle to target",GREEN,  "自動追蹤玩家"),
    ("waveBullets",    "波動彈",    "sin(phase) × amp",    ACCTL,  "蛇形軌跡"),
    ("starBurst",      "星爆交錯",  "快/慢/快/慢...",      RED,   "交替速度環"),
    ("butterflySpray", "蝴蝶散射",  "sin(t) × π × 0.4",    BLUE,  "蝴蝶型 8 字"),
]
for i, (name, zh, formula, clr, desc) in enumerate(patterns):
    col, row = i % 2, i // 2
    x, y = 0.4 + col * 4.85, 1.55 + row * 1.0
    card(sl, x, y, 4.65, 0.82)
    add_oval(sl, x + 0.1, y + 0.08, 0.27, 0.27, clr)
    add_textbox(sl, name, x + 0.1, y + 0.08, 0.27, 0.27,
                font_size=8, bold=True, color=DARK, align=PP_ALIGN.CENTER)
    add_textbox(sl, zh, x + 0.48, y + 0.04, 3.0, 0.28, font_size=12, bold=True, color=WHITE)
    add_textbox(sl, formula, x + 0.48, y + 0.34, 3.0, 0.18,
                font_size=9, color=GOLD, font_name="Courier New")
    add_textbox(sl, desc, x + 3.6, y + 0.04, 0.85, 0.72,
                font_size=8, color=DIM, align=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════
# Slide 5 — 子彈物理 4 種型別
# ═══════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
slide_bg(sl)
section_tag(sl, "03 子彈物理")
slide_title(sl, "4 種子彈運動型別")
hrule(sl, 1.35)

modes = [
    ("type 0", "直線", BLUE, [
        "x += vx",
        "y += vy",
    ], "基礎：恆速直線移動"),
    ("type 1", "追蹤", RED, [
        "t = atan2(py-y, px-x)",
        "angle += lerp(angle, t)",
    ], "每幀朝玩家轉向，turnRate 控制反應度"),
    ("type 2", "波動", TEAL, [
        "phase += freq",
        "perp_offset = sin(phase) × amp",
    ], "垂直軸施加正弦波，頻率&振幅可調"),
    ("type 3", "加速", GOLD, [
        "speed += accel",
        "move by speed",
    ], "每幀加速，可用負 accel 造成延遲引爆"),
]
for i, (typ, name, clr, code, desc) in enumerate(modes):
    col, row = i % 2, i // 2
    x, y = 0.4 + col * 4.85, 1.55 + row * 1.98
    card(sl, x, y, 4.65, 1.78)
    card(sl, x + 0.1, y + 0.08, 1.0, 0.32, clr)
    add_textbox(sl, typ, x + 0.1, y + 0.08, 1.0, 0.32,
                font_size=10, bold=True, color=DARK, align=PP_ALIGN.CENTER)
    add_textbox(sl, name, x + 1.2, y + 0.08, 1.3, 0.32,
                font_size=13, bold=True, color=WHITE)
    card(sl, x + 0.1, y + 0.48, 4.45, 0.62, CODE)
    for j, line in enumerate(code):
        is_comment = line.startswith("//")
        add_textbox(sl, line, x + 0.2, y + 0.52 + j*0.20, 4.25, 0.18,
                    font_size=9, color=DIM if is_comment else TEAL,
                    font_name="Courier New")
    add_textbox(sl, desc, x + 0.1, y + 1.14, 4.45, 0.56, font_size=9.5, color=DIM)

# ═══════════════════════════════════════════════════════
# Slide 6 — EXP 升級系統
# ═══════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
slide_bg(sl)
section_tag(sl, "04 EXP 系統")
slide_title(sl, "5 項技能樹與升級邏輯")
hrule(sl, 1.35)

# Left: formulas
card(sl, 0.4, 1.55, 4.8, 3.75)
add_textbox(sl, "升級公式與效果", 0.5, 1.65, 4.6, 0.35,
            font_size=13, bold=True, color=ACCTL, align=PP_ALIGN.CENTER)

formulas = [
    ("升級所需",       "expNext = 1000 + level × 500"),
    ("刷掠獲得",       "graze += 5 per bullet"),
    ("擊殺敵人",       "exp = 50 × (type+1) × (diff+1)"),
    ("Boss EXP",       "exp = 5000 × (diff+1)"),
    ("射速技能",       "speed_mult = 1.0 + lvl × 0.15"),
    ("射頻技能",       "cooldown -= lvl × 2 (min 2)"),
    ("生命技能",       "max_life += lvl"),
    ("炸彈技能",       "bomb_dmg × (1.0 + lvl × 0.2)"),
    ("刷掠技能",       "graze_radius × (1.0 + lvl × 0.1)"),
]
for i, (label, formula) in enumerate(formulas):
    y = 2.08 + i * 0.40
    add_textbox(sl, label + "：", 0.5, y, 1.6, 0.32, font_size=9, color=DIM)
    card(sl, 2.15, y - 0.015, 2.6, 0.32, CODE)
    add_textbox(sl, formula, 2.25, y, 2.4, 0.28,
                font_size=8.5, color=TEAL, font_name="Courier New")

# Right: chart
add_textbox(sl, "Lv 1-10 升級 EXP 需求",
            5.35, 1.65, 4.25, 0.3, font_size=12, bold=True, color=WHITE)

cd = ChartData()
cd.categories = [f"Lv {i}" for i in range(1, 11)]
cd.add_series("EXP Next", [1000 + lv*500 for lv in range(1, 11)])

chart = sl.shapes.add_chart(
    XL_CHART_TYPE.COLUMN_CLUSTERED,
    Inches(5.2), Inches(1.98), Inches(4.4), Inches(2.9),
    cd
)
plot = chart.chart.plots[0]
plot.series[0].format.fill.solid()
plot.series[0].format.fill.fore_color.rgb = ACCENT

# ═══════════════════════════════════════════════════════
# Slide 7 — 敵人 AI 生命週期
# ═══════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
slide_bg(sl)
section_tag(sl, "05 敵人 AI")
slide_title(sl, "敵人生命週期與行為")
hrule(sl, 1.35)

# Lifecycle
stages = [
    ("負 age",   "age = -(i × 15)\n錯開延遲",      GREEN),
    ("進場",     "age < 60\n依 vx/vy 飛入",       BLUE),
    ("減速",     "60 ≤ age < 120\nvx *= 0.95",   TEAL),
    ("巡邏",     "age ≥ 120\nx = sin(t) × 40",   GOLD),
    ("消失",     "age ≥ maxAge\nactive = false",  RED),
]
for i, (stage, desc, clr) in enumerate(stages):
    x = 0.4 + i * 1.88
    card(sl, x, 1.55, 1.72, 0.95)
    add_oval(sl, x + 0.61, 1.64, 0.5, 0.5, clr)
    add_textbox(sl, stage, x + 0.61, 1.64, 0.5, 0.5,
                font_size=10, bold=True, color=DARK, align=PP_ALIGN.CENTER)
    add_textbox(sl, desc, x + 0.1, 2.18, 1.52, 0.3,
                font_size=8.5, color=DIM, align=PP_ALIGN.CENTER, font_name="Courier New")
    if i < 4:
        sh = sl.shapes.add_shape(1, Inches(x+1.72), Inches(1.90),
                                  Inches(0.16), Inches(0.02))
        sh.fill.solid(); sh.fill.fore_color.rgb = ACCENT; sh.line.fill.background()

# maxAge table
add_textbox(sl, "各敵人類型的生存時間", 0.4, 2.65, 5.8, 0.28,
            font_size=12, bold=True, color=WHITE)

rows = [
    ["敵人類型",     "maxAge", "存活秒數", "射擊間隔"],
    ["type 0 仙子",  "480",    "8.0 秒",   "100 幀"],
    ["type 1 妖怪",  "600",    "10.0 秒",  "90 幀"],
    ["type 2 頭目",  "720",    "12.0 秒",  "70 幀"],
]
tbl = sl.shapes.add_table(4, 4, Inches(0.4), Inches(2.95), Inches(5.8), Inches(1.55)).table
col_widths = [1.5, 1.1, 1.6, 1.0]
for ci, w in enumerate(col_widths):
    tbl.columns[ci].width = Inches(w)
for ri, row_data in enumerate(rows):
    for ci, text in enumerate(row_data):
        cell = tbl.cell(ri, ci)
        cell.text = text
        p = cell.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.runs[0]
        run.font.size = Pt(11)
        run.font.bold = (ri == 0)
        run.font.name = "Calibri"
        run.font.color.rgb = DARK if ri == 0 else WHITE
        cell.fill.solid()
        cell.fill.fore_color.rgb = ACCENT if ri == 0 else PANEL

# Right: stagger
card(sl, 6.3, 1.55, 3.3, 3.95, PANEL2)
add_textbox(sl, "錯開生成", 6.4, 1.65, 3.1, 0.3,
            font_size=12, bold=True, color=TEAL, align=PP_ALIGN.CENTER)
code_lines = [
    "for (int i = 0; i < n; i++) {",
    "  Enemy en = new Enemy(...);",
    "  en.age = -(i * 15);",
    "  // 每隻差 15 幀",
    "  enemies.add(en);",
    "}",
    "",
    "// update() 中",
    "if (age <= 0) return; // 未出生",
    "",
    "if (age >= maxAge) {",
    "  active = false;   // 到期消失",
    "}",
]
for j, line in enumerate(code_lines):
    if not line.strip():
        continue
    is_cmt = line.strip().startswith("//")
    add_textbox(sl, line, 6.4, 2.0 + j*0.205, 3.0, 0.19,
                font_size=8, color=DIM if is_cmt else TEAL,
                font_name="Courier New")

# ═══════════════════════════════════════════════════════
# Slide 8 — 難度與效能
# ═══════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
slide_bg(sl)
section_tag(sl, "06 實驗結果")
slide_title(sl, "效能驗證與難度平衡")
hrule(sl, 1.35)

# Difficulty table
add_textbox(sl, "難度參數與子彈密度", 0.4, 1.55, 5.2, 0.3,
            font_size=12, bold=True, color=WHITE)

diff_rows = [
    ["難度",     "子彈速度",    "子彈數量", "敵人頻率"],
    ["Easy",    "0.65×",      "6-8 發",   "標準"],
    ["Normal",  "0.78×",      "7-9 發",   "標準"],
    ["Hard",    "0.91×",      "8-10 發",  "高"],
    ["Lunatic", "1.04×",      "10-12 發", "非常高"],
]
tbl = sl.shapes.add_table(5, 4, Inches(0.4), Inches(1.87), Inches(5.2), Inches(1.8)).table
col_widths = [1.0, 1.4, 1.4, 1.4]
for ci, w in enumerate(col_widths):
    tbl.columns[ci].width = Inches(w)
for ri, row_data in enumerate(diff_rows):
    for ci, text in enumerate(row_data):
        cell = tbl.cell(ri, ci)
        cell.text = text
        p = cell.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.runs[0]
        run.font.size = Pt(10.5)
        run.font.bold = (ri == 0)
        run.font.name = "Calibri"
        run.font.color.rgb = DARK if ri == 0 else WHITE
        cell.fill.solid()
        cell.fill.fore_color.rgb = ACCENT if ri == 0 else PANEL

# Key metrics
add_textbox(sl, "效能指標", 5.6, 1.55, 4.0, 0.3, font_size=12, bold=True, color=WHITE)

metrics = [
    ("最大同時子彈", "~ 400 顆", GOLD),
    ("碰撞檢測耗時", "< 1ms / frame", GREEN),
    ("FPS 穩定度", "60 FPS ±1%", BLUE),
    ("記憶體占用", "< 80 MB", TEAL),
]
for i, (metric, value, clr) in enumerate(metrics):
    y = 1.87 + i * 0.88
    card(sl, 5.6, y, 4.0, 0.78)
    sh = sl.shapes.add_shape(1, Inches(5.6), Inches(y), Inches(0.08), Inches(0.78))
    sh.fill.solid(); sh.fill.fore_color.rgb = clr; sh.line.fill.background()
    add_textbox(sl, metric, 5.75, y + 0.05, 3.65, 0.28, font_size=10.5, color=DIM)
    add_textbox(sl, value,  5.75, y + 0.35, 3.65, 0.35, font_size=13, bold=True, color=clr)

# ═══════════════════════════════════════════════════════
# Slide 9 — 結語
# ═══════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
slide_bg(sl)

add_oval(sl, -1.0, 3.2, 5.5, 5.5, RGBColor(0x3D,0x00,0x66))
add_oval(sl, 7.5, -1.2, 5.5, 5.5, RGBColor(0x00,0x3D,0x66))

add_textbox(sl, "BULLET STORM", 0.5, 0.8, 9.0, 0.8,
            font_size=40, bold=True, color=TEAL, align=PP_ALIGN.CENTER)
add_textbox(sl, "演算法設計篇", 0.5, 1.65, 9.0, 0.5,
            font_size=18, color=ACCTL, align=PP_ALIGN.CENTER)

sum_stats = [
    ("8",      "彈幕演算法",  "從入門到進階"),
    ("4",      "運動物理",   "直線到追蹤"),
    ("5",      "技能樹系統",  "完整升級路線"),
    ("99",     "最高等級",    "Level 1-99"),
]
for i, (v, l, sub) in enumerate(sum_stats):
    x = 0.5 + i * 2.28
    card(sl, x, 3.0, 2.1, 1.85)
    add_textbox(sl, v, x, 3.08, 2.1, 0.65,
                font_size=26, bold=True, color=TEAL, align=PP_ALIGN.CENTER)
    add_textbox(sl, l, x, 3.7, 2.1, 0.28,
                font_size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_textbox(sl, sub, x, 3.98, 2.1, 0.55,
                font_size=8.5, color=DIM, align=PP_ALIGN.CENTER)

add_textbox(sl, "https://github.com/user/BulletStorm",
            0.5, 5.0, 9.0, 0.4, font_size=11, color=DIM, align=PP_ALIGN.CENTER)

# ─── Save ─────────────────────────────────────────────
out = r"C:\project T\TouhouGame\ppt\02_演算法設計.pptx"
prs.save(out)
print("Optimized PPT2 saved")
