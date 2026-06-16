# -*- coding: utf-8 -*-
"""
複製 01_圖形介面設計.pptx → 01_圖形介面設計_含遊戲例圖.pptx
並插入 2 張遊戲實機截圖投影片（不更動原 8 張）：
  - 原 index 1 後：「遊戲實機畫面：視窗佈局」（gameplay + char/diff select）
  - 原 index 4 後：「HUD 特寫與選角畫面」（hud + charselect 並排）
"""
import shutil, os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

PPT_DIR = r"C:\project T\BulletStorm\ppt"
SRC = os.path.join(PPT_DIR, "01_圖形介面設計.pptx")
DST = os.path.join(PPT_DIR, "01_圖形介面設計_含遊戲例圖.pptx")
ASSET = os.path.join(PPT_DIR, "assets")

# 沿用 rebuild_ppt1.py 的配色
BG    = RGBColor(0x0A, 0x0A, 0x18)
PANEL = RGBColor(0x10, 0x15, 0x2B)
CYAN  = RGBColor(0x00, 0xC8, 0xD7)
CYANL = RGBColor(0x40, 0xE0, 0xFF)
GOLD  = RGBColor(0xFF, 0xD7, 0x00)
WHITE = RGBColor(0xEA, 0xF0, 0xF8)
DIM   = RGBColor(0x80, 0x98, 0xB0)
GREEN = RGBColor(0x4D, 0xFF, 0x8C)
BLUE  = RGBColor(0x00, 0xBB, 0xFF)
PURP  = RGBColor(0xC7, 0x7D, 0xFF)

shutil.copyfile(SRC, DST)
prs = Presentation(DST)

def bg(sl): sl.background.fill.solid(); sl.background.fill.fore_color.rgb = BG

def box(sl, text, x, y, w, h, sz=14, bold=False, color=WHITE,
        align=PP_ALIGN.LEFT, italic=False, font="Calibri"):
    tb = sl.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = text
    r.font.size = Pt(sz); r.font.bold = bold; r.font.italic = italic
    r.font.color.rgb = color; r.font.name = font

def rect(sl, x, y, w, h, fill=PANEL, line=None):
    sh = sl.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    sh.fill.solid(); sh.fill.fore_color.rgb = fill
    if line: sh.line.color.rgb = line; sh.line.width = Pt(0.8)
    else: sh.line.fill.background()
    return sh

def oval(sl, x, y, w, h, fill):
    sh = sl.shapes.add_shape(9, Inches(x), Inches(y), Inches(w), Inches(h))
    sh.fill.solid(); sh.fill.fore_color.rgb = fill
    sh.line.fill.background()

def tag(sl, text):
    rect(sl, 0.4, 0.15, 2.6, 0.38, fill=CYAN)
    box(sl, text, 0.4, 0.15, 2.6, 0.38, sz=11, bold=True, color=BG, align=PP_ALIGN.CENTER)

def title(sl, text):
    box(sl, text, 0.4, 0.62, 9.2, 0.72, sz=26, bold=True, color=WHITE)

def blank():
    return prs.slides.add_slide(prs.slide_layouts[6])

# ════════════════════════════════════════════════════════
# 新投影片 A — 遊戲實機畫面：完整視窗佈局
# 插入位置：「系統架構：視窗佈局與座標」(index 1) 之後
# ════════════════════════════════════════════════════════
slA = blank(); bg(slA)
tag(slA, "02  System Architecture")
title(slA, "遊戲實機畫面：完整視窗佈局")

box(slA, "左側 480×560 遊戲戰場 + 右側 HUD 面板，以 Java Swing 雙重緩衝即時渲染。",
    0.4, 1.28, 9.2, 0.32, sz=12, color=CYANL)

# 主圖：gameplay（16:12 ≈ 800×600）
iw = 7.0; ih = iw * 600 / 800
slA.shapes.add_picture(
    os.path.join(ASSET, "gui_01_gameplay.png"),
    Inches(0.4), Inches(1.68), Inches(iw), Inches(ih)
)

# 右側兩張小圖（char select + diff select）
sw = 2.35; sh_img = sw * 600 / 800
sx = 0.4 + iw + 0.18
slA.shapes.add_picture(
    os.path.join(ASSET, "gui_03_charselect.png"),
    Inches(sx), Inches(1.68), Inches(sw), Inches(sh_img)
)
box(slA, "角色選擇", sx, 1.68 + sh_img + 0.04, sw, 0.26,
    sz=10, color=CYANL, align=PP_ALIGN.CENTER)

slA.shapes.add_picture(
    os.path.join(ASSET, "gui_04_diffselect.png"),
    Inches(sx), Inches(1.68 + sh_img + 0.34), Inches(sw), Inches(sh_img)
)
box(slA, "難度選擇", sx, 1.68 + sh_img * 2 + 0.38, sw, 0.26,
    sz=10, color=CYANL, align=PP_ALIGN.CENTER)

box(slA,
    "遊戲區邊框：綠色 (32,16) 480×560  ·  HUD 面板邊框：藍色 x=544  ·  背景：150 顆動態星星",
    0.4, 5.3, 9.2, 0.28, sz=10, color=DIM, italic=True)

# ════════════════════════════════════════════════════════
# 新投影片 B — HUD 特寫與選角畫面
# 插入位置：「HUD 設計、角色與難度系統」(原 index 4，現在是 index 5) 之後
# ════════════════════════════════════════════════════════
slB = blank(); bg(slB)
tag(slB, "03  Technologies Used")
title(slB, "遊戲實機畫面：HUD 特寫與角色系統")

box(slB, "左：HUD 右側面板放大（全 8 個元素）；右：角色選擇實機截圖（Aurora Striker 選中狀態）。",
    0.4, 1.28, 9.2, 0.32, sz=12, color=CYANL)

# Left: HUD closeup (square 560×560 → show as 4.5" square)
hw = 4.4; hh = hw
slB.shapes.add_picture(
    os.path.join(ASSET, "gui_02_hud.png"),
    Inches(0.4), Inches(1.68), Inches(hw), Inches(hh)
)

# Labels for HUD elements (right of image)
lx = 0.4 + hw + 0.22
items = [
    (GOLD,  "HIGH SCORE / SCORE", "金色分數顯示，玩家最重要資訊"),
    (GREEN, "HEALTH 生命值",       "綠色愛心×3，剩餘生命一目了然"),
    (CYANL, "BOMBS 炸彈",         "橘色圖示，庫存即時更新"),
    (PURP,  "LV + EXP 條",        "等級 07、紫色升級進度條 68%"),
    (BLUE,  "SKILLS ×5 技能",     "5 種技能各顯示 0–5 顆星"),
]
for i, (c, h, d) in enumerate(items):
    y = 1.72 + i * 0.78
    oval(slB, lx, y + 0.06, 0.18, 0.18, c)
    box(slB, h, lx + 0.26, y + 0.02, 4.7, 0.28, sz=11, bold=True, color=WHITE)
    box(slB, d, lx + 0.26, y + 0.28, 4.7, 0.26, sz=10, color=DIM)

# Right: char select (16:12 → fit in ~4.5"×3.4")
cw2 = 4.4; ch2 = cw2 * 600 / 800
slB.shapes.add_picture(
    os.path.join(ASSET, "gui_03_charselect.png"),
    Inches(lx), Inches(1.68 + len(items)*0.78 + 0.08),
    Inches(cw2), Inches(ch2)
)
box(slB, "角色選擇：Aurora Striker（金框選中）vs Chronos Engineer",
    lx, 1.68 + len(items)*0.78 + ch2 + 0.12, cw2, 0.28, sz=10, color=DIM, italic=True)

# ════════════════════════════════════════════════════════
# 重排投影片順序（原 8 張不動，新 A B 插入到正確位置）
# ════════════════════════════════════════════════════════
sldIdLst = prs.slides._sldIdLst
ids = list(sldIdLst)
# 目前：原 0-7（共 8 張）+ 新 A=8, B=9
a_id, b_id = ids[-2], ids[-1]
sldIdLst.remove(a_id)
sldIdLst.remove(b_id)
# 原 index 1 = 系統架構：視窗佈局 → A 插到 index 2
sldIdLst.insert(2, a_id)
# 原 index 4 = HUD設計（現在因為 A 插入，變成 index 5）→ B 插到 index 6
sldIdLst.insert(6, b_id)

prs.save(DST)
print("Saved:", DST, "| total slides =", len(prs.slides._sldIdLst))
