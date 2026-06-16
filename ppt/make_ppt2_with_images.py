# -*- coding: utf-8 -*-
"""
在【不更動原本 9 張投影片】的前提下，複製 02_演算法設計.pptx 為新檔，
並於「8 種彈幕演算法」那張之後，插入 2 張以遊戲實機渲染的例圖投影片。
原檔 02_演算法設計.pptx 不受影響。
"""
import shutil, os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

PPT_DIR = r"C:\project T\BulletStorm\ppt"
SRC = os.path.join(PPT_DIR, "02_演算法設計.pptx")
DST = os.path.join(PPT_DIR, "02_演算法設計_含遊戲例圖.pptx")
ASSET = os.path.join(PPT_DIR, "assets")

# —— 沿用 rebuild_ppt2.py 的配色 ——
BG    = RGBColor(0x0F,0x0A,0x1E)
PANEL = RGBColor(0x1A,0x12,0x30)
PURP  = RGBColor(0x9B,0x5D,0xE5)
PURPL = RGBColor(0xC7,0x7D,0xFF)
TEAL  = RGBColor(0x00,0xF5,0xD4)
GOLD  = RGBColor(0xFE,0xE4,0x40)
WHITE = RGBColor(0xF0,0xEA,0xF8)
DIM   = RGBColor(0x8B,0x80,0xA8)
RED   = RGBColor(0xF1,0x5B,0xB5)
GREEN = RGBColor(0x57,0xCC,0x99)

shutil.copyfile(SRC, DST)
prs = Presentation(DST)

def bg(sl): sl.background.fill.solid(); sl.background.fill.fore_color.rgb = BG

def box(sl,text,x,y,w,h,sz=14,bold=False,color=WHITE,align=PP_ALIGN.LEFT,
        italic=False,font="Calibri"):
    tb=sl.shapes.add_textbox(Inches(x),Inches(y),Inches(w),Inches(h))
    tf=tb.text_frame; tf.word_wrap=True
    p=tf.paragraphs[0]; p.alignment=align
    r=p.add_run(); r.text=text
    r.font.size=Pt(sz); r.font.bold=bold; r.font.italic=italic
    r.font.color.rgb=color; r.font.name=font

def rect(sl,x,y,w,h,fill=PANEL,line=None):
    sh=sl.shapes.add_shape(1,Inches(x),Inches(y),Inches(w),Inches(h))
    sh.fill.solid(); sh.fill.fore_color.rgb=fill
    if line: sh.line.color.rgb=line; sh.line.width=Pt(1)
    else: sh.line.fill.background()
    return sh

def oval(sl,x,y,w,h,fill):
    sh=sl.shapes.add_shape(9,Inches(x),Inches(y),Inches(w),Inches(h))
    sh.fill.solid(); sh.fill.fore_color.rgb=fill; sh.line.fill.background()

def tag(sl,text):
    rect(sl,0.4,0.15,3.0,0.38,fill=PURP)
    box(sl,text,0.4,0.15,3.0,0.38,sz=11,bold=True,color=WHITE,align=PP_ALIGN.CENTER)

def title(sl,text):
    box(sl,text,0.4,0.62,9.2,0.72,sz=26,bold=True,color=WHITE)

def blank():
    return prs.slides.add_slide(prs.slide_layouts[6])

# ════════════════════════════════════════════════════════
# 新投影片 A — 8 種彈幕「遊戲實機畫面」總覽
# ════════════════════════════════════════════════════════
slA = blank(); bg(slA)
tag(slA,"03  Technologies Used")
title(slA,"演算法實機畫面：8 種彈幕")
box(slA,"以下皆由遊戲原始 BulletPattern 演算法即時渲染，套用遊戲相同的發光彈芯繪製。",
    0.4,1.28,9.2,0.32,sz=12,color=PURPL)
# gallery image (1160×584 ≈ 1.986:1)
gw=9.2; gh=gw*584/1160
slA.shapes.add_picture(os.path.join(ASSET,"00_gallery.png"),
                       Inches(0.4),Inches(1.7),Inches(gw),Inches(gh))
box(slA,"環形爆散 · 瞄準扇射 · 雙環交轉 · 螺旋臂　|　追蹤彈 · 波動彈 · 星爆交錯 · 蝴蝶散射",
    0.4,1.7+gh+0.12,9.2,0.32,sz=11,color=DIM,italic=True,align=PP_ALIGN.CENTER)

# ════════════════════════════════════════════════════════
# 新投影片 B — 重點演算法特寫 + 公式對照
# ════════════════════════════════════════════════════════
slB = blank(); bg(slB)
tag(slB,"03  Technologies Used")
title(slB,"重點演算法特寫")

feats=[
    ("01_circularBurst.png","環形爆散","θᵢ = θ₀ + i×(2π/n)","n 發等角灑成整圈",TEAL),
    ("04_spiralArm.png",    "螺旋臂",  "offset += Δθ / frame","逐幀偏移→旋轉螺旋",GREEN),
    ("05_homingBullets.png","追蹤彈",  "angle = lerp(a, target, 0.05)","每幀朝玩家修正",RED),
]
cw=2.96
for i,(img,zh,formula,desc,c) in enumerate(feats):
    x=0.4+i*3.07
    rect(slB,x,1.5,cw,3.7)
    # image (square 560×560), shrunk & centred so caption text stays inside the card
    iw=2.4
    ix=x+(cw-iw)/2
    slB.shapes.add_picture(os.path.join(ASSET,img),
                           Inches(ix),Inches(1.62),Inches(iw),Inches(iw))
    yb=1.62+iw+0.1
    oval(slB,x+0.14,yb+0.03,0.2,0.2,c)
    box(slB,zh,x+0.42,yb-0.03,cw-0.5,0.3,sz=14,bold=True,color=WHITE)
    box(slB,formula,x+0.16,yb+0.32,cw-0.32,0.3,sz=10,color=GOLD,font="Courier New")
    box(slB,desc,x+0.16,yb+0.62,cw-0.32,0.3,sz=10.5,color=DIM)

box(slB,"螺旋臂為連續多幀疊加結果；追蹤彈呈現朝右下角玩家收斂的彎曲軌跡。",
    0.4,5.34,9.2,0.28,sz=10,color=DIM,italic=True)

# ════════════════════════════════════════════════════════
# 把兩張新投影片移到「8 種彈幕生成演算法」(原 S4, index 3) 之後
# 原本各張順序與內容完全不變
# ════════════════════════════════════════════════════════
sldIdLst = prs.slides._sldIdLst
ids = list(sldIdLst)               # 目前順序：原 0-8，新 A=9, B=10
a_id, b_id = ids[-2], ids[-1]
sldIdLst.remove(a_id); sldIdLst.remove(b_id)
# 原 S4 是 index 3 → 插在 index 4、5
sldIdLst.insert(4, a_id)
sldIdLst.insert(5, b_id)

prs.save(DST)
print("Saved:", DST, "| total slides =", len(prs.slides._sldIdLst))
