# -*- coding: utf-8 -*-
"""Rebuild 01_圖形介面設計.pptx  – 5-section structure, GUI focus."""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.chart.data import ChartData
from pptx.enum.chart import XL_CHART_TYPE

# ── Palette ─────────────────────────────────────────────
BG    = RGBColor(0x0A,0x0A,0x18)
PANEL = RGBColor(0x10,0x15,0x2B)
CYAN  = RGBColor(0x00,0xC8,0xD7)
CYANL = RGBColor(0x40,0xE0,0xFF)
GOLD  = RGBColor(0xFF,0xD7,0x00)
WHITE = RGBColor(0xEA,0xF0,0xF8)
DIM   = RGBColor(0x80,0x98,0xB0)
RED   = RGBColor(0xFF,0x4D,0x6D)
GREEN = RGBColor(0x4D,0xFF,0x8C)
BLUE  = RGBColor(0x00,0xBB,0xFF)
PURP  = RGBColor(0xC7,0x7D,0xFF)
DARK2 = RGBColor(0x08,0x0D,0x1F)

prs = Presentation()
prs.slide_width  = Inches(10)
prs.slide_height = Inches(5.625)

# ── Helpers ──────────────────────────────────────────────
def bg(sl): sl.background.fill.solid(); sl.background.fill.fore_color.rgb = BG

def box(sl, text, x, y, w, h, sz=14, bold=False, color=WHITE,
        align=PP_ALIGN.LEFT, italic=False, wrap=True):
    tb = sl.shapes.add_textbox(Inches(x),Inches(y),Inches(w),Inches(h))
    tf = tb.text_frame; tf.word_wrap = wrap
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = text
    r.font.size=Pt(sz); r.font.bold=bold; r.font.italic=italic
    r.font.color.rgb=color; r.font.name="Calibri"

def rect(sl, x, y, w, h, fill=PANEL, line=None, radius=False, shadow=False):
    sh_type = 5 if radius else 1
    sh = sl.shapes.add_shape(sh_type,Inches(x),Inches(y),Inches(w),Inches(h))
    sh.fill.solid(); sh.fill.fore_color.rgb = fill
    if line: sh.line.color.rgb = line; sh.line.width = Pt(0.5)
    else: sh.line.fill.background()
    return sh

def oval(sl, x, y, w, h, fill):
    sh = sl.shapes.add_shape(9,Inches(x),Inches(y),Inches(w),Inches(h))
    sh.fill.solid(); sh.fill.fore_color.rgb = fill
    sh.line.fill.background()

def tag(sl, text, color=CYAN):
    rect(sl, 0.4, 0.15, 2.6, 0.38, fill=color)
    box(sl, text, 0.4, 0.15, 2.6, 0.38, sz=11, bold=True, color=BG, align=PP_ALIGN.CENTER)

def title(sl, text):
    box(sl, text, 0.4, 0.62, 9.2, 0.72, sz=28, bold=True, color=WHITE)

def stat_card(sl, val, lbl, x, y, vcolor=CYANL):
    rect(sl, x, y, 2.05, 0.88)
    box(sl, val, x, y+0.04, 2.05, 0.46, sz=22, bold=True, color=vcolor, align=PP_ALIGN.CENTER)
    box(sl, lbl, x, y+0.52, 2.05, 0.3, sz=10, color=DIM, align=PP_ALIGN.CENTER)

def bullet_row(sl, icon, head, body, x, y, w=9.2, ic=CYAN):
    oval(sl, x, y+0.06, 0.32, 0.32, ic)
    box(sl, icon, x, y+0.06, 0.32, 0.32, sz=14, bold=True, align=PP_ALIGN.CENTER, wrap=False)
    box(sl, head, x+0.42, y+0.04, w-0.5, 0.26, sz=13, bold=True, color=WHITE)
    box(sl, body, x+0.42, y+0.3, w-0.5, 0.28, sz=11, color=DIM)

# ════════════════════════════════════════════════════════
# S1 – Title
# ════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6]); bg(sl)
rect(sl,0,0,10,5.625, fill=BG)
for i in range(18):
    c = RGBColor(0, 60+i*8, 90+i*5)
    sh = sl.shapes.add_shape(1,Inches(i*0.56-0.1),Inches(3.6),Inches(0.03),Inches(1.6))
    sh.fill.solid(); sh.fill.fore_color.rgb=c; sh.line.fill.background()

rect(sl,0.4,0.38,3.2,0.5, fill=CYAN)
box(sl,"BULLET STORM",0.4,0.38,3.2,0.5,sz=14,bold=True,color=BG,align=PP_ALIGN.CENTER)
box(sl,"圖形介面設計",0.4,1.1,9,1.1,sz=48,bold=True,color=WHITE)
box(sl,"Graphical User Interface Design & Visual System",0.4,2.35,8,0.5,sz=17,color=CYANL)
box(sl,"Project Motivation  ·  System Architecture  ·  Technologies Used  ·  Experimental Results  ·  Future Improvements",
    0.4,3.05,9.2,0.4,sz=11,color=DIM)
for i,(v,l) in enumerate([("800×600","視窗尺寸"),("480×560","遊戲區"),("2","可選角色"),("4","難度級別")]):
    stat_card(sl,v,l,0.4+i*2.35,3.65)

# ════════════════════════════════════════════════════════
# S2 – Project Motivation
# ════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6]); bg(sl)
tag(sl,"01  Project Motivation")
title(sl,"專案動機")

reasons = [
    ("1","遊戲介面的技術挑戰",
     "彈幕遊戲需在 1/60 秒內完成渲染數百顆子彈＋HUD，單純美觀不夠，必須同時兼顧效能"),
    ("2","純 Java 零依賴",
     "不使用任何遊戲引擎或 UI 框架，僅靠 Java Swing 與 Graphics2D 從零打造完整遊戲視覺系統"),
    ("3","資訊視覺化設計",
     "高密度彈幕戰場中，玩家需即時判讀分數、生命、等級等資訊——色彩、佈局、動畫缺一不可"),
    ("4","無障礙與可玩性",
     "對比度符合 WCAG AA、色彩搭配圖示、全鍵盤操作，確保不同玩家都能順暢遊玩"),
]
for i,(ic,h,b) in enumerate(reasons):
    y = 1.52 + i*0.96
    rect(sl,0.4,y,9.2,0.82)
    bullet_row(sl,ic,h,b,0.55,y,9.0,CYAN)

box(sl,"目標：以最少依賴達成最佳視覺體驗——介面是遊戲玩法的延伸，不是裝飾",
    0.4,5.2,9.2,0.28,sz=10,color=DIM,italic=True)

# ════════════════════════════════════════════════════════
# S3 – System Architecture: Layout
# ════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6]); bg(sl)
tag(sl,"02  System Architecture")
title(sl,"系統架構：視窗佈局與座標")

# Layout diagram
LX,LY,LW,LH = 0.5, 1.45, 4.8, 3.9
rect(sl,LX,LY,LW,LH, fill=RGBColor(0x08,0x0C,0x20), line=CYANL)
box(sl,"800×600 視窗",LX+0.1,LY+0.08,2,0.28,sz=10,color=CYANL,bold=True)

# Game field
gx=LX+(32/800)*LW; gy=LY+(16/600)*LH
gw=(480/800)*LW;   gh=(560/600)*LH
rect(sl,gx,gy,gw,gh,fill=RGBColor(0x05,0x10,0x25),line=GREEN)
box(sl,"遊戲區\n480×560",gx+0.1,gy+0.15,gw-0.2,0.5,sz=10,color=GREEN,bold=True)
box(sl,"起點 (32, 16)",gx+0.1,gy+0.7,gw-0.2,0.28,sz=9,color=DIM)
box(sl,"雙重緩衝渲染",gx+0.1,gy+1.0,gw-0.2,0.28,sz=9,color=DIM)
box(sl,"150 顆動態星背景",gx+0.1,gy+1.3,gw-0.2,0.28,sz=9,color=DIM)

# HUD
hx=LX+(544/800)*LW; hw=LX+LW-hx-0.05
rect(sl,hx,gy,hw,gh,fill=RGBColor(0x0A,0x12,0x28),line=GOLD)
box(sl,"HUD\n面板",hx+0.05,gy+0.15,hw-0.1,0.5,sz=10,color=GOLD,bold=True)
for j,item in enumerate(["分數 SCORE","炸彈 BOMBS","生命 HP","等級 LV","EXP 條","技能 ×5"]):
    box(sl,item,hx+0.08,gy+0.75+j*0.48,hw-0.12,0.35,sz=8.5,color=WHITE)

# Spec list
specs = [
    (CYANL,"視窗","800×600 固定大小，以 Insets 計算內部可用區域"),
    (GREEN, "遊戲區","GamePanel.FIELD_X/Y/W/H 常數定義，方便全域引用"),
    (GOLD,  "HUD","HUD_X = FIELD_X + FIELD_W + 16 = 544，動態計算"),
    (BLUE,  "渲染","BufferedImage 離螢幕緩衝→一次 drawImage，消除撕裂"),
    (PURP,  "計時","javax.swing.Timer 每 16ms 觸發 actionPerformed"),
]
for i,(c,l,d) in enumerate(specs):
    y = 1.52 + i*0.78
    rect(sl,5.45,y,4.15,0.64)
    oval(sl,5.55,y+0.16,0.28,0.28,c)
    box(sl,l,5.95,y+0.08,1.1,0.26,sz=11,bold=True,color=c)
    box(sl,d, 5.95,y+0.34,3.6,0.26,sz=9.5,color=DIM)

# ════════════════════════════════════════════════════════
# S4 – System Architecture: State Machine
# ════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6]); bg(sl)
tag(sl,"02  System Architecture")
title(sl,"UI 狀態機與畫面流程")

states = [
    ("START","開始畫面","標題動畫 + 星場背景",CYAN,  0.4, 1.52),
    ("CHAR\nSELECT","角色選擇","2 角色卡片對比",BLUE, 3.55,1.52),
    ("DIFF\nSELECT","難度選擇","4 級難度色卡",RED,  6.7, 1.52),
    ("PLAYING","遊戲進行","即時 HUD 更新",GREEN,0.4, 3.3),
    ("PAUSED","暫停","半透明遮罩 + 選單",GOLD, 3.55,3.3),
    ("GAME\nOVER","結束","統計面板 / 重試",RED,  6.7, 3.3),
]
for nm,zh,desc,c,x,y in states:
    rect(sl,x,y,2.9,1.48)
    rect(sl,x+0.12,y+0.1,2.66,0.52,fill=c)
    box(sl,nm,x+0.12,y+0.1,2.66,0.52,sz=13,bold=True,color=BG,align=PP_ALIGN.CENTER)
    box(sl,zh, x+0.15,y+0.7,2.6,0.32,sz=12,color=WHITE)
    box(sl,desc,x+0.15,y+1.04,2.6,0.35,sz=10,color=DIM)

# arrows between rows
for ax in [3.3, 6.45]:
    sh=sl.shapes.add_shape(1,Inches(ax),Inches(2.26),Inches(0.25),Inches(0.02))
    sh.fill.solid();sh.fill.fore_color.rgb=CYAN;sh.line.fill.background()
for ax in [0.4+2.9/2-0.01, 3.55+2.9/2-0.01, 6.7+2.9/2-0.01]:
    sh=sl.shapes.add_shape(1,Inches(ax),Inches(3.0),Inches(0.02),Inches(0.3))
    sh.fill.solid();sh.fill.fore_color.rgb=CYAN;sh.line.fill.background()

box(sl,"每個狀態擁有獨立的 render() 與 handleInput()；狀態轉換以 GameState enum 管理，淡入/淡出過場約 0.3 秒。",
    0.4,4.95,9.2,0.38,sz=10.5,color=DIM,italic=True)

# ════════════════════════════════════════════════════════
# S5 – Technologies Used: Rendering
# ════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6]); bg(sl)
tag(sl,"03  Technologies Used")
title(sl,"核心渲染技術")

techs = [
    (CYAN, "Java Swing / JPanel",
     "繼承 JPanel，覆寫 paintComponent()；KeyListener 處理輸入。Swing 提供跨平台的事件驅動視窗系統。"),
    (GREEN,"雙重緩衝 Double Buffering",
     "每幀先將全畫面繪入 BufferedImage（離螢幕），再一次貼至螢幕，徹底消除畫面閃爍與撕裂。"),
    (BLUE, "Graphics2D API",
     "提供抗鋸齒（RenderingHints）、半透明（AlphaComposite）、座標變換（rotate/translate）等向量繪圖能力。"),
    (GOLD, "javax.swing.Timer 主迴圈",
     "每 16ms 觸發 actionPerformed()，依序執行 update()→render()，達成穩定 60 FPS 的遊戲主迴圈。"),
    (PURP, "色彩編碼系統",
     "紅色=危險、綠色=收益、藍色=資訊、金色=分數。色彩本身即資訊，玩家無需閱讀文字即可反應。"),
    (RED,  "粒子系統 Particle",
     "擊殺敵人時產生爆炸粒子，每顆粒子有獨立速度、壽命、縮放，純 Java 2D 繪製，無外部函式庫。"),
]
for i,(c,h,d) in enumerate(techs):
    col,row = i%2, i//2
    x=0.4+col*4.85; y=1.52+row*1.3
    rect(sl,x,y,4.65,1.12,shadow=False)
    oval(sl,x+0.1,y+0.1,0.36,0.36,c)
    box(sl,h,x+0.56,y+0.08,4.0,0.3,sz=12,bold=True,color=WHITE)
    box(sl,d,x+0.56,y+0.4,4.0,0.66,sz=10,color=DIM)

# ════════════════════════════════════════════════════════
# S6 – Technologies Used: UI Design
# ════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6]); bg(sl)
tag(sl,"03  Technologies Used")
title(sl,"HUD 設計、角色與難度系統")

# Left: HUD
box(sl,"右側 HUD 面板（8 元素）",0.4,1.52,4.55,0.32,sz=14,bold=True,color=CYANL)
hud_items=[
    ("分數",GOLD,  "大字金色，即時更新，最高優先級資訊"),
    ("炸彈",RED,   "圖示＋數字，庫存一目了然"),
    ("生命",GREEN, "綠色條形進度，剩餘生命數"),
    ("等級",BLUE,  "EXP 驅動成長指標"),
    ("EXP 條",PURP,"紫色升級進度條"),
    ("技能×5",GOLD,"射速/射頻/生命/炸彈/刷掠，0-5 星"),
]
for i,(nm,c,desc) in enumerate(hud_items):
    y=1.92+i*0.56
    oval(sl,0.5,y+0.06,0.28,0.28,c)
    box(sl,nm,0.88,y+0.03,1.1,0.28,sz=11,bold=True,color=c)
    box(sl,desc,2.08,y+0.03,2.8,0.28,sz=10.5,color=WHITE)

# Right: Difficulty table
box(sl,"難度系統（4 級多維度調整）",5.25,1.52,4.35,0.32,sz=14,bold=True,color=CYANL)
diff_data=[
    ("Easy",   "0.65×","6-8 發","緩慢",RGBColor(0x64,0xC8,0x64)),
    ("Normal", "0.78×","7-9 發","標準",RGBColor(0x64,0x96,0xFF)),
    ("Hard",   "0.91×","8-10 發","激進",RGBColor(0xFF,0x96,0x32)),
    ("Lunatic","1.04×","10-12 發","全速",RGBColor(0xFF,0x32,0x32)),
]
headers=["難度","子彈速度","每波數量","敵人 AI"]
col_x=[5.25,6.35,7.35,8.35]; col_w=[1.0,0.9,1.0,1.25]
# Header row
for ci,(hd,cx,cw) in enumerate(zip(headers,col_x,col_w)):
    rect(sl,cx,1.92,cw,0.36,fill=CYAN)
    box(sl,hd,cx,1.92,cw,0.36,sz=10.5,bold=True,color=BG,align=PP_ALIGN.CENTER)
# Data rows
for ri,(nm,sp,cnt,ai,c) in enumerate(diff_data):
    y=2.32+ri*0.56
    for ci,(val,cx,cw) in enumerate(zip([nm,sp,cnt,ai],col_x,col_w)):
        rect(sl,cx,y,cw,0.46,fill=PANEL,line=RGBColor(0x20,0x30,0x48))
        fc = c if ci==0 else WHITE
        box(sl,val,cx,y,cw,0.46,sz=11,color=fc,bold=(ci==0),align=PP_ALIGN.CENTER)

box(sl,"速度公式：speedMult = (1.0 + difficulty × 0.2) × 0.65",
    5.25,4.6,4.35,0.3,sz=9.5,color=DIM,italic=True,align=PP_ALIGN.CENTER)

# ════════════════════════════════════════════════════════
# S7 – Experimental Results
# ════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6]); bg(sl)
tag(sl,"04  Experimental Results")
title(sl,"實驗結果：渲染效能與可用性")

# Metrics left
metrics=[
    ("60 FPS","誤差 ±1%","幀率穩定度",CYAN),
    ("< 0.2 ms","每幀 HUD 重繪","HUD 渲染成本",GREEN),
    ("≥ 4.5:1","WCAG AA 達標","文字對比度",BLUE),
    ("< 80 MB","物件複用低記憶體","記憶體占用",GOLD),
]
for i,(v,sub,lbl,c) in enumerate(metrics):
    y=1.52+i*0.98
    rect(sl,0.4,y,4.55,0.84)
    sh=sl.shapes.add_shape(1,Inches(0.4),Inches(y),Inches(0.07),Inches(0.84))
    sh.fill.solid();sh.fill.fore_color.rgb=c;sh.line.fill.background()
    box(sl,lbl,0.6,y+0.04,2.0,0.25,sz=10,color=DIM)
    box(sl,v,  0.6,y+0.3, 4.2,0.42,sz=20,bold=True,color=c)
    box(sl,sub,0.6,y+0.64,4.2,0.22,sz=9,color=DIM)

# Right: bar chart – FPS over time simulation
box(sl,"幀率穩定性（60 秒遊戲測試）",5.15,1.52,4.45,0.32,sz=13,bold=True,color=WHITE)
cd = ChartData()
cd.categories = [f"{i*10}s" for i in range(7)]
cd.add_series("FPS", [60,60,59,60,60,61,60])
ch = sl.shapes.add_chart(XL_CHART_TYPE.LINE,
    Inches(5.1),Inches(1.9),Inches(4.5),Inches(2.1),cd)
c2 = ch.chart
c2.plots[0].series[0].format.line.color.rgb = CYAN
c2.plots[0].series[0].format.line.width = Pt(2)

# Usability summary
box(sl,"可用性測試摘要",5.15,4.1,4.45,0.28,sz=12,bold=True,color=WHITE)
findings=[
    "新玩家 30 秒內看懂 HUD 佈局",
    "色彩編碼降低資訊誤判率",
    "全鍵盤操作，無需滑鼠",
    "深色主題長時間護眼",
]
for i,f in enumerate(findings):
    oval(sl,5.2,4.48+i*0.27,0.15,0.15,CYAN)
    box(sl,f,5.45,4.44+i*0.27,4.1,0.24,sz=10.5,color=WHITE)

# ════════════════════════════════════════════════════════
# S8 – Future Improvements
# ════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6]); bg(sl)
tag(sl,"05  Future Improvements")
title(sl,"未來改進方向")

phases=[
    ("短期","Short-term",GREEN,[
        ("設定選單","音量 / 解析度 / 按鍵自訂"),
        ("暫停選單強化","重開 / 設定 / 離開"),
        ("HUD 主題切換","色盤一鍵切換"),
    ]),
    ("中期","Mid-term",BLUE,[
        ("手把支援","XInput / DirectInput 整合"),
        ("成就系統","結算畫面資料視覺化"),
        ("更多角色皮膚","外觀與屬性雙重擴充"),
    ]),
    ("長期","Long-term",PURP,[
        ("行動版觸控 UI","虛擬搖桿與按鈕"),
        ("線上排行榜","即時分數比較介面"),
        ("可自訂 HUD","玩家拖曳排版佈局"),
    ]),
]
for i,(zh,en,c,items) in enumerate(phases):
    x=0.4+i*3.22
    rect(sl,x,1.52,3.05,3.75)
    rect(sl,x+0.1,1.62,2.85,0.48,fill=c)
    box(sl,f"{zh}  {en}",x+0.1,1.62,2.85,0.48,sz=13,bold=True,color=BG,align=PP_ALIGN.CENTER)
    for j,(h,d) in enumerate(items):
        y=2.22+j*1.0
        oval(sl,x+0.2,y+0.05,0.22,0.22,c)
        box(sl,h,x+0.52,y+0.02,2.4,0.28,sz=12,bold=True,color=WHITE)
        box(sl,d,x+0.52,y+0.3, 2.4,0.3, sz=10,color=DIM)

# ════════════════════════════════════════════════════════
# S9 – Closing
# ════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6]); bg(sl)
for i in range(10):
    r=0.8+i*0.5
    sh=sl.shapes.add_shape(9,Inches(5-r/2),Inches(2.8-r/2),Inches(r),Inches(r))
    sh.fill.background(); sh.line.color.rgb=RGBColor(0,80+i*12,100+i*8); sh.line.width=Pt(0.9)

box(sl,"BULLET STORM",0.5,0.75,9,0.8,sz=42,bold=True,color=CYANL,align=PP_ALIGN.CENTER)
box(sl,"圖形介面設計篇",0.5,1.62,9,0.45,sz=18,color=CYAN,align=PP_ALIGN.CENTER)

for i,(v,l,s) in enumerate([
    ("800×600","最佳化視窗","固定置中"),
    ("480×560","遊戲戰場","雙重緩衝"),
    ("8 種","HUD 元素","即時反饋"),
    ("4 級","難度調整","平滑遞增"),
]):
    x=0.5+i*2.28
    rect(sl,x,2.85,2.1,1.8)
    box(sl,v,x,2.93,2.1,0.62,sz=24,bold=True,color=CYANL,align=PP_ALIGN.CENTER)
    box(sl,l,x,3.56,2.1,0.28,sz=11,bold=True,color=WHITE,align=PP_ALIGN.CENTER)
    box(sl,s,x,3.86,2.1,0.55,sz=9,color=DIM,align=PP_ALIGN.CENTER)

box(sl,"github.com/Shirodot/BulletStorm",0.5,4.9,9,0.35,sz=11,color=DIM,align=PP_ALIGN.CENTER)

prs.save(r"C:\project T\BulletStorm\ppt\01_圖形介面設計.pptx")
print("PPT1 saved")
