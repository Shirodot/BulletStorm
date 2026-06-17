# -*- coding: utf-8 -*-
"""
就地校正《02_演算法設計_含遊戲例圖.pptx》，使其與 src/ 程式碼 100% 對齊，
並補上遺漏內容。原則：只做精準文字替換與底部註腳新增，不更動既有版面。
本腳本為 idempotent（可重複執行，不會重複新增註腳）。

對應修正（皆已逐項比對 src/）：
  1. Boss 難度公式：Boss 用 ×0.25（Boss.java:93），非投影片所列的 ×0.20（那是一般敵人）。
  2. type 3（加速）已實作但未在遊戲中啟用 → 加註腳。
  3. homing 實為夾限轉向 clamp(±0.05)（Bullet.java:61-63），非 lerp。
  4. 巡邏含垂直分量、減速階段 vx/vy 同乘 0.95（Enemy.java:54-64）→ 加註腳。
  5. 兩位具名 Boss：Commander Aurora / Engineer Chronos（Boss.java:53）。
  6. Boss 命中 EXP：100 ×(diff+1)（GamePanel.java:224）。
  7. 5 種已實作未啟用彈幕（BulletPattern.java）。
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

PPT_DIR = os.path.dirname(os.path.abspath(__file__))
FILE = os.path.join(PPT_DIR, "02_演算法設計_含遊戲例圖.pptx")

# —— 沿用 make_ppt2_with_images.py 的配色 ——
DIM   = RGBColor(0x8B, 0x80, 0xA8)
GOLD  = RGBColor(0xFE, 0xE4, 0x40)
WHITE = RGBColor(0xF0, 0xEA, 0xF8)

prs = Presentation(FILE)
slides = list(prs.slides)


def slide_text(sl):
    return "\n".join(sh.text_frame.text for sh in sl.shapes if sh.has_text_frame)


def replace_substr(sl, old, new, guard=None):
    """在 slide 內所有 run 將 old 子字串替換為 new。
    guard 若已存在於該頁文字中則跳過（確保 append 類修改 idempotent）。"""
    if guard and guard in slide_text(sl):
        return 0
    hits = 0
    for sh in sl.shapes:
        if not sh.has_text_frame:
            continue
        for p in sh.text_frame.paragraphs:
            for r in p.runs:
                if old in r.text:
                    r.text = r.text.replace(old, new)
                    hits += 1
    return hits


def add_footnote(sl, x, y, w, h, text, guard, sz=9, color=DIM, italic=True,
                 align=PP_ALIGN.LEFT):
    """於指定位置加一個小字註腳；若該頁已含 guard 則跳過。"""
    if guard in slide_text(sl):
        return 0
    tb = sl.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.size = Pt(sz)
    r.font.italic = italic
    r.font.color.rgb = color
    r.font.name = "Calibri"
    return 1


# ── 文字替換 (slide_index, old, new, guard) ─────────────────────────────
text_ops = [
    # 5. 兩位具名 Boss（架構頁註腳尾端 append）
    (2, "達成演算法複用與解耦。",
        "達成演算法複用與解耦。　兩位具名 Boss：Commander Aurora、Engineer Chronos。",
        "Commander Aurora"),
    # 3. homing 夾限轉向（彈幕公式頁）
    (3, "angle = lerp(angle, target, 0.05)",
        "angle += clamp(target - angle, ±0.05)", None),
    # 3. homing 夾限轉向（重點特寫頁）
    (5, "angle = lerp(a, target, 0.05)",
        "angle += clamp(target - a, ±0.05)", None),
    # 6. Boss 命中 EXP
    (6, "5000 × (diff+1)",
        "5000×(diff+1) 擊殺／100×(diff+1) 命中", "命中"),
    # 4. 減速階段 vx、vy 同乘 0.95
    (7, "vx *= 0.95", "vx,vy *= 0.95", None),
]

for idx, old, new, guard in text_ops:
    n = replace_substr(slides[idx], old, new, guard)
    print(f"  slide{idx}: replace '{old[:24]}...' -> hits={n}")

# ── 底部註腳新增 ────────────────────────────────────────────────────────
# 7. 5 種已實作未啟用彈幕（彈幕公式頁）
add_footnote(
    slides[3], 0.40, 5.30, 9.20, 0.28,
    "＊以上 8 種為遊戲實際啟用；另有 aimFan / spinningRing / accelBurst / "
    "laserBarrage / sakuraPetal 共 5 種已實作、暫未啟用（BulletPattern 共 13 種）。",
    guard="暫未啟用", sz=9)

# 2. type 3 已實作保留擴充（子彈物理頁，置左下 type 3 區下方）
add_footnote(
    slides[6], 0.40, 5.40, 4.70, 0.20,
    "＊type 0–2 於遊戲即時運作；type 3（加速，可為負 accel）已實作，目前保留作擴充。",
    guard="保留作擴充", sz=8.5)

# 4. 巡邏垂直分量 + 減速 vy（敵人 AI 頁，置表格下方左側）
add_footnote(
    slides[7], 0.40, 5.30, 5.90, 0.24,
    "＊巡邏：x = entryX + sin(t)×40，y = entryY+60 + sin(t·½)×15；"
    "減速階段 vx、vy 同乘 0.95。",
    guard="sin(t·½)", sz=9)

# ── 1. Boss 難度公式（實驗結果頁）：重寫為兩條，明確區分一般敵人 / Boss ──
DIFF_SHAPE_OLD = "難度公式"
for sh in slides[8].shapes:
    if sh.has_text_frame and DIFF_SHAPE_OLD in sh.text_frame.text:
        tf = sh.text_frame
        tf.word_wrap = True
        # 放大容器以容納兩行
        sh.top = Inches(4.55)
        sh.height = Inches(0.95)
        tf.clear()
        lines = [
            "難度公式　一般敵人　speedMult = (1.0 + diff×0.2)×0.65 → "
            "Easy 0.65 / Normal 0.78 / Hard 0.91 / Lunatic 1.04",
            "　　　　　　Boss　　speedMult = (1.0 + diff×0.25)×0.65 → "
            "0.65 / 0.81 / 0.97 / 1.14",
        ]
        for i, line in enumerate(lines):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            r = p.add_run()
            r.text = line
            r.font.size = Pt(10)
            r.font.italic = True
            r.font.color.rgb = GOLD if i == 1 else DIM
            r.font.name = "Calibri"
        print("  slide8: difficulty formula rewritten (enemy 0.2 / boss 0.25)")
        break

prs.save(FILE)
print("Saved:", FILE, "| slides =", len(prs.slides._sldIdLst))
