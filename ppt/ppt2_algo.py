"""
PPT 2 — Bullet Storm：演算法設計篇
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.chart.data import ChartData
from pptx.enum.chart import XL_CHART_TYPE
import math

# ─── Palette (Purple/Math theme) ────────────────────────
DARK   = RGBColor(0x0F, 0x0A, 0x1E)
PANEL  = RGBColor(0x1A, 0x12, 0x30)
ACCENT = RGBColor(0x9B, 0x5D, 0xE5)
ACCTL  = RGBColor(0xC7, 0x7D, 0xFF)
TEAL   = RGBColor(0x00, 0xF5, 0xD4)
GOLD   = RGBColor(0xFE, 0xE4, 0x40)
WHITE  = RGBColor(0xF0, 0xEA, 0xF8)
DIM    = RGBColor(0x8B, 0x80, 0xA8)
RED    = RGBColor(0xF1, 0x5B, 0xB5)
BLUE   = RGBColor(0x00, 0xBB, 0xF9)
GREEN  = RGBColor(0x57, 0xCC, 0x99)
PANEL2 = RGBColor(0x13, 0x0D, 0x22)
CODE   = RGBColor(0x0A, 0x07, 0x18)

prs = Presentation()
prs.slide_width  = Inches(10)
prs.slide_height = Inches(5.625)

# ─── Helpers ────────────────────────────────────────────
def slide_bg(slide):
    bg = slide.background; f = bg.fill
    f.solid(); f.fore_color.rgb = DARK

def card(slide, x, y, w, h, fill=None):
    c = fill or PANEL
    sh = slide.shapes.add_shape(5, Inches(x), Inches(y), Inches(w), Inches(h))
    sh.fill.solid(); sh.fill.fore_color.rgb = c
    sh.line.color.rgb = RGBColor(0x2A,0x1E,0x42); sh.line.width = Pt(0.5)
    return sh

def add_textbox(slide, text, x, y, w, h,
                font_size=14, bold=False, color=WHITE,
                align=PP_ALIGN.LEFT, italic=False, font_name="Calibri"):
    txb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = txb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = align
    run = p.add_run()
    run.text = text; run.font.size = Pt(font_size)
    run.font.bold = bold; run.font.italic = italic
    run.font.color.rgb = color; run.font.name = font_name
    return txb

def add_oval(slide, x, y, w, h, fill_rgb):
    sh = slide.shapes.add_shape(9, Inches(x), Inches(y), Inches(w), Inches(h))
    sh.fill.solid(); sh.fill.fore_color.rgb = fill_rgb
    sh.line.fill.background(); return sh

def section_tag(slide, text):
    card(slide, 0.4, 0.15, 2.3, 0.4, ACCENT)
    add_textbox(slide, text, 0.4, 0.15, 2.3, 0.4,
                font_size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

def slide_title(slide, text):
    add_textbox(slide, text, 0.4, 0.6, 9.2, 0.7, font_size=28, bold=True, color=WHITE)

def hrule(slide, y):
    sh = slide.shapes.add_shape(1, Inches(0.4), Inches(y), Inches(9.2), Inches(0.03))
    sh.fill.solid(); sh.fill.fore_color.rgb = ACCENT; sh.line.fill.background()

# ═══════════════════════════════════════════════════════
# Slide 1 — Title
# ═══════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
slide_bg(sl)

# Orbit rings (ovals with no fill)
for r, clr, t in [(7, ACCENT, 70), (5, TEAL, 75)]:
    sh = sl.shapes.add_shape(9, Inches(5.5 - r/2), Inches(-r/2), Inches(r), Inches(r))
    sh.fill.background()
    sh.line.color.rgb = clr; sh.line.width = Pt(1.2)

# Badge
card(sl, 0.5, 0.45, 2.8, 0.52, ACCENT)
add_textbox(sl, "BULLET STORM", 0.5, 0.45, 2.8, 0.52,
            font_size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

add_textbox(sl, "演算法設計", 0.5, 1.15, 7, 1.2, font_size=46, bold=True, color=WHITE)
add_textbox(sl, "Algorithm Design & Internal Mechanics",
            0.5, 2.45, 7, 0.55, font_size=18, color=ACCTL)

sh = sl.shapes.add_shape(1, Inches(0.5), Inches(3.1), Inches(3.5), Inches(0.04))
sh.fill.solid(); sh.fill.fore_color.rgb = TEAL; sh.line.fill.background()

add_textbox(sl, "12 種彈幕演算法  ·  EXP 升級系統  ·  敵人 AI\n碰撞偵測  ·  物件生命週期  ·  難度自適應",
            0.5, 3.22, 7, 0.85, font_size=14, color=DIM)

stats = [("12","彈幕演算法"), ("5","技能樹系統"), ("4","子彈運動型"), ("O(n²)","碰撞複雜度")]
for i, (v, l) in enumerate(stats):
    x = 0.5 + i * 2.3
    card(sl, x, 4.35, 2.1, 0.95)
    add_textbox(sl, v, x, 4.38, 2.1, 0.48, font_size=22, bold=True, color=TEAL, align=PP_ALIGN.CENTER)
    add_textbox(sl, l, x, 4.84, 2.1, 0.3,  font_size=10, color=DIM,  align=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════
# Slide 2 — 目錄
# ═══════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
slide_bg(sl)
section_tag(sl, "CONTENTS")
slide_title(sl, "簡報目錄")
hrule(sl, 1.35)

toc = [
    ("01","專案動機",   "彈幕遊戲的演算法挑戰性與設計動機"),
    ("02","系統架構",   "類別設計、物件關係、資料流向"),
    ("03","使用技術",   "12 種彈幕演算法 + 子彈運動模式詳解"),
    ("04","實驗結果",   "效能分析、EXP 驗證、難度曲線測試"),
    ("05","未來改進",   "機器學習難度、複雜彈幕公式"),
]
for i, (n, t, d) in enumerate(toc):
    y = 1.52 + i * 0.8
    card(sl, 0.4, y, 9.2, 0.65)
    add_oval(sl, 0.55, y + 0.1, 0.45, 0.45, ACCENT)
    add_textbox(sl, n, 0.55, y + 0.1, 0.45, 0.45,
                font_size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_textbox(sl, t, 1.12, y + 0.07, 2.5, 0.3, font_size=14, bold=True, color=WHITE)
    add_textbox(sl, d, 1.12, y + 0.35, 7.9, 0.25, font_size=11, color=DIM)

# ═══════════════════════════════════════════════════════
# Slide 3 — 專案動機
# ═══════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
slide_bg(sl)
section_tag(sl, "01 專案動機")
slide_title(sl, "演算法挑戰與設計動機")
hrule(sl, 1.35)

challenges = [
    ("🔢", "數學密集型運算",   "三角函數、向量運算、參數方程式在每幀執行數百次"),
    ("🎯", "碰撞偵測效率",    "子彈×敵人×玩家的 O(n²) 計算，需空間分割優化"),
    ("🤖", "敵人行為決策",    "自動入場路徑、巡邏 AI、射擊時機判斷皆需演算法"),
    ("⚖", "難度動態平衡",    "4 種難度 × 速度/數量/頻率的多維度調整公式"),
]
for i, (icon, title, desc) in enumerate(challenges):
    y = 1.55 + i * 0.93
    card(sl, 0.4, y, 5.5, 0.78)
    add_textbox(sl, icon, 0.52, y + 0.14, 0.55, 0.5, font_size=22, align=PP_ALIGN.CENTER)
    add_textbox(sl, title, 1.15, y + 0.07, 4.6, 0.3,  font_size=13, bold=True, color=WHITE)
    add_textbox(sl, desc,  1.15, y + 0.38, 4.6, 0.32, font_size=10, color=DIM)

card(sl, 6.1, 1.55, 3.5, 3.7, PANEL2)
add_textbox(sl, "設計原則", 6.2, 1.62, 3.3, 0.38,
            font_size=14, bold=True, color=ACCTL, align=PP_ALIGN.CENTER)
principles = [
    "數學美感優先（視覺優美軌跡）",
    "可組合性：12 個獨立函式",
    "參數化設計：難度即係數",
    "物件複用：Bullet 通用型別",
    "延遲清除：避免迭代中修改",
    "隨機性 + 確定性混合使用",
]
for i, p in enumerate(principles):
    add_oval(sl, 6.25, 2.12 + i*0.5, 0.17, 0.17, TEAL)
    add_textbox(sl, p, 6.52, 2.06 + i*0.5, 2.85, 0.3, font_size=11, color=WHITE)

# ═══════════════════════════════════════════════════════
# Slide 4 — 系統架構：類別設計
# ═══════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
slide_bg(sl)
section_tag(sl, "02 系統架構")
slide_title(sl, "核心類別架構設計")
hrule(sl, 1.35)

classes = [
    ("GamePanel",        "遊戲主控制器",          ACCENT, 3.9, 1.55),
    ("BulletPattern",    "彈幕演算法庫 (static)", TEAL,   0.4, 2.75),
    ("Boss",             "Boss 射擊邏輯",         RED,    3.55, 2.75),
    ("Enemy",            "敵人 AI + 射擊",        BLUE,   6.8,  2.75),
    ("Bullet",           "子彈物理 (4型)",         GOLD,   0.4,  3.9),
    ("Player",           "玩家控制 + EXP",        GREEN,  3.55, 3.9),
    ("ExperienceSystem", "升級/技能計算",         ACCTL,  6.8,  3.9),
]
for (name, role, clr, x, y) in classes:
    card(sl, x, y, 2.5, 0.72)
    card(sl, x + 0.08, y + 0.06, 2.34, 0.3, clr)
    add_textbox(sl, name, x + 0.08, y + 0.06, 2.34, 0.3,
                font_size=11, bold=True, color=DARK, align=PP_ALIGN.CENTER)
    add_textbox(sl, role, x + 0.1, y + 0.38, 2.3, 0.28,
                font_size=10, color=WHITE, align=PP_ALIGN.CENTER)

# Connecting lines
lines = [
    (5.15, 2.25, 1.65, 2.75), (5.15, 2.25, 4.8, 2.75), (5.15, 2.25, 8.05, 2.75),
    (1.65, 3.47, 1.65, 3.9),  (4.8,  3.47, 4.8, 3.9),  (8.05, 3.47, 8.05, 3.9),
]
for (x1, y1, x2, y2) in lines:
    sh = sl.shapes.add_shape(1,
        Inches(min(x1,x2)), Inches(min(y1,y2)),
        Inches(abs(x2-x1)) or Inches(0.02), Inches(abs(y2-y1)) or Inches(0.02))
    sh.fill.solid(); sh.fill.fore_color.rgb = ACCENT; sh.line.fill.background()

add_textbox(sl, "─ ─  呼叫/使用關係   │   BulletPattern 為純靜態工具類，Boss/Enemy 呼叫其方法產生子彈 List",
            0.4, 4.88, 9.2, 0.28, font_size=10, color=DIM)

# ═══════════════════════════════════════════════════════
# Slide 5 — 12 種彈幕演算法
# ═══════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
slide_bg(sl)
section_tag(sl, "03 使用技術")
slide_title(sl, "12 種彈幕模式演算法")
hrule(sl, 1.35)

patterns = [
    ("circularBurst",  "環形爆散",  "angle = offset + i×(2π/n)",    ACCTL),
    ("spiralArm",      "螺旋臂",    "offset += Δθ per frame",         TEAL),
    ("aimedSpread",    "瞄準扇射",  "base = atan2(dy,dx) ± spread",  BLUE),
    ("doubleRing",     "雙環",      "outer > 0，inner < 0 速度",     GOLD),
    ("starBurst",      "星爆交錯",  "i%2==0→fast，i%2==1→slow",     RED),
    ("homingBullets",  "追蹤彈",    "angle lerp to target, t=0.05",  GREEN),
    ("waveBullets",    "波動彈",    "x+=vx + perp×sin(phase)×amp",  ACCTL),
    ("spinningRing",   "旋轉環",    "angle += angularVel/frame",     TEAL),
    ("accelBurst",     "加速彈",    "speed += accel per frame",       BLUE),
    ("butterflySpray", "蝴蝶散射",  "spread = sin(t)×π×0.4",        GOLD),
    ("laserBarrage",   "雷射彈幕",  "angle = target ± rand×spread", RED),
    ("sakuraPetal",    "花瓣",      "petal×spread + speed lerp",    GREEN),
]
for i, (name, zh, formula, clr) in enumerate(patterns):
    col, row = i % 3, i // 3
    x, y = 0.4 + col * 3.22, 1.55 + row * 1.02
    card(sl, x, y, 3.05, 0.88)
    add_textbox(sl, zh + "  " + name + "()",
                x + 0.1, y + 0.05, 2.85, 0.3, font_size=11, bold=True, color=WHITE)
    add_oval(sl, x + 0.1, y + 0.44, 0.2, 0.2, clr)
    add_textbox(sl, formula, x + 0.38, y + 0.53, 2.6, 0.28,
                font_size=9, color=DIM, font_name="Courier New")

# ═══════════════════════════════════════════════════════
# Slide 6 — 子彈物理 4 種型別
# ═══════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
slide_bg(sl)
section_tag(sl, "03 使用技術")
slide_title(sl, "子彈物理運動模式（4 型）")
hrule(sl, 1.35)

modes = [
    ("type 0", "標準直線", BLUE, [
        "x += vx",
        "y += vy",
        "// angularVel 旋轉彈道",
    ], "最基礎模式。速度向量每幀累加\n座標，可加 angularVel 形成螺旋。"),
    ("type 1", "追蹤 Homing", RED, [
        "t = atan2(py-y,px-x)",
        "d = normalize(t-angle)",
        "angle += clamp(d,±0.05)",
    ], "每幀計算玩家方向，以 turnRate\n漸進轉向，高壓迫感追蹤彈。"),
    ("type 2", "波動 Wave", TEAL, [
        "phase += frequency",
        "perpX = -sin(angle)",
        "x += vx+perpX×sin(p)×a",
    ], "垂直軸施加正弦波偏移。\nfreq 控頻率，amp 控幅度。"),
    ("type 3", "加速 Accel", GOLD, [
        "speed += accel",
        "vx = cos(angle)×speed",
        "vy = sin(angle)×speed",
    ], "每幀增加速度。可初始負速\n製造「先退後衝」的突進效果。"),
]
for i, (typ, name, clr, code, desc) in enumerate(modes):
    col, row = i % 2, i // 2
    x, y = 0.4 + col * 4.85, 1.55 + row * 2.02
    card(sl, x, y, 4.65, 1.82)
    card(sl, x + 0.1, y + 0.08, 1.0, 0.32, clr)
    add_textbox(sl, typ, x + 0.1, y + 0.08, 1.0, 0.32,
                font_size=10, bold=True, color=DARK, align=PP_ALIGN.CENTER)
    add_textbox(sl, name, x + 1.2, y + 0.08, 3.3, 0.32,
                font_size=14, bold=True, color=WHITE)
    # Code block
    card(sl, x + 0.1, y + 0.48, 4.45, 0.8, CODE)
    for j, line in enumerate(code):
        is_comment = line.startswith("//")
        add_textbox(sl, line, x + 0.2, y + 0.52 + j*0.24, 4.25, 0.22,
                    font_size=9.5, color=DIM if is_comment else TEAL,
                    font_name="Courier New")
    add_textbox(sl, desc, x + 0.1, y + 1.32, 4.45, 0.44, font_size=10, color=DIM)

# ═══════════════════════════════════════════════════════
# Slide 7 — EXP 升級系統
# ═══════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
slide_bg(sl)
section_tag(sl, "03 使用技術")
slide_title(sl, "EXP 升級系統演算法")
hrule(sl, 1.35)

# Left: formulas
card(sl, 0.4, 1.55, 4.45, 3.75)
add_textbox(sl, "升級公式與技能係數", 0.5, 1.65, 4.25, 0.38,
            font_size=14, bold=True, color=ACCTL, align=PP_ALIGN.CENTER)

formulas = [
    ("升級所需 EXP",   "exp_next = 1000 + level × 500"),
    ("刷掠 EXP",       "graze_exp = 5 / bullet"),
    ("擊殺 EXP",       "kill_exp = 50 × (type+1) × (diff+1)"),
    ("Boss EXP",       "boss_exp = 5000 × (diff+1)"),
    ("射速技能",       "speed_mult = 1.0 + lvl × 0.15"),
    ("射頻技能",       "cooldown -= lvl × 2   (min = 2)"),
    ("生命技能",       "lives += hp_level"),
    ("炸彈技能",       "bomb_dmg × (1.0 + lvl × 0.2)"),
]
for i, (label, formula) in enumerate(formulas):
    y = 2.12 + i * 0.46
    add_textbox(sl, label + "：", 0.5, y, 1.8, 0.35, font_size=10, color=DIM)
    card(sl, 2.3, y - 0.02, 2.4, 0.35, CODE)
    add_textbox(sl, formula, 2.38, y, 2.28, 0.3,
                font_size=9, color=TEAL, font_name="Courier New")

# Right: EXP bar chart
add_textbox(sl, "升級 EXP 成長曲線（Level 1-10）",
            5.05, 1.65, 4.55, 0.35, font_size=13, bold=True, color=WHITE)

cd = ChartData()
cd.categories = [f"Lv {i}" for i in range(1, 11)]
cd.add_series("升級所需 EXP", [1000 + lv*500 for lv in range(1, 11)])

chart = sl.shapes.add_chart(
    XL_CHART_TYPE.COLUMN_CLUSTERED,
    Inches(4.9), Inches(2.0), Inches(4.8), Inches(3.3),
    cd
)
chart.chart.plots[0].series[0].format.fill.solid()
chart.chart.plots[0].series[0].format.fill.fore_color.rgb = ACCENT

# ═══════════════════════════════════════════════════════
# Slide 8 — 敵人 AI 生命週期
# ═══════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
slide_bg(sl)
section_tag(sl, "03 使用技術")
slide_title(sl, "敵人 AI 行為與生命週期管理")
hrule(sl, 1.35)

# Lifecycle stages
stages = [
    ("生成",  "age = -i×15\n錯開入場延遲",           GREEN),
    ("進場",  "age < 60\n依 vx/vy 飛入場地",          BLUE),
    ("減速",  "60 < age < 120\nvx, vy × 0.95",       TEAL),
    ("巡邏",  "age ≥ 120\nx = sin(t)×40",            GOLD),
    ("消失",  "age ≥ maxAge\nactive = false",          RED),
]
for i, (stage, desc, clr) in enumerate(stages):
    x = 0.4 + i * 1.88
    card(sl, x, 1.55, 1.72, 1.0, PANEL)
    add_oval(sl, x + 0.61, 1.65, 0.5, 0.5, clr)
    add_textbox(sl, stage, x + 0.61, 1.65, 0.5, 0.5,
                font_size=11, bold=True, color=DARK, align=PP_ALIGN.CENTER)
    add_textbox(sl, desc, x + 0.1, 2.2, 1.52, 0.32,
                font_size=9, color=DIM, align=PP_ALIGN.CENTER, font_name="Courier New")
    if i < 4:
        arr = sl.shapes.add_shape(1, Inches(x+1.72), Inches(2.04),
                                   Inches(0.16), Inches(0.02))
        arr.fill.solid(); arr.fill.fore_color.rgb = ACCENT; arr.line.fill.background()

# maxAge table
add_textbox(sl, "各敵人類型存活時間 (maxAge)", 0.4, 2.75, 5.5, 0.32,
            font_size=13, bold=True, color=WHITE)

rows = [
    ["敵人類型",     "maxAge (幀)", "存活秒數 @60fps", "射擊間隔"],
    ["仙子 (type 0)", "480",        "8.0 秒",          "100 幀"],
    ["妖怪 (type 1)", "600",        "10.0 秒",         "90 幀"],
    ["頭目 (type 2)", "720",        "12.0 秒",         "70 幀"],
]
from pptx.util import Inches as In
tbl = sl.shapes.add_table(4, 4, In(0.4), In(3.12), In(5.5), In(1.65)).table
col_widths = [1.5, 1.3, 1.8, 1.0]
for ci, w in enumerate(col_widths):
    tbl.columns[ci].width = In(w)
for ri, row in enumerate(rows):
    for ci, text in enumerate(row):
        cell = tbl.cell(ri, ci)
        cell.text = text
        p = cell.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.runs[0]
        run.font.size = Pt(12)
        run.font.bold = (ri == 0)
        run.font.name = "Calibri"
        run.font.color.rgb = DARK if ri == 0 else WHITE
        cell.fill.solid()
        cell.fill.fore_color.rgb = ACCENT if ri == 0 else PANEL

# Right: stagger code
card(sl, 6.1, 1.55, 3.5, 3.75, PANEL2)
add_textbox(sl, "錯開生成演算法", 6.2, 1.65, 3.3, 0.35,
            font_size=13, bold=True, color=TEAL, align=PP_ALIGN.CENTER)
code_lines = [
    "// 生成波次時設定負 age",
    "for (int i = 0; i < n; i++) {",
    "  Enemy en = new Enemy(...);",
    "  en.age = -(i * 15);",
    "  // 每隻差 15 幀入場",
    "  enemies.add(en);",
    "}",
    "",
    "// update() 跳過未出生的",
    "if (age <= 0) return;",
    "",
    "// maxAge 到達後自動消失",
    "if (age >= maxAge) {",
    "  active = false; return;",
    "}",
]
for j, line in enumerate(code_lines):
    if not line.strip():
        continue
    is_cmt = line.strip().startswith("//")
    add_textbox(sl, line, 6.2, 2.08 + j*0.22, 3.2, 0.2,
                font_size=8.5, color=DIM if is_cmt else TEAL,
                font_name="Courier New")

# ═══════════════════════════════════════════════════════
# Slide 9 — 實驗結果
# ═══════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
slide_bg(sl)
section_tag(sl, "04 實驗結果")
slide_title(sl, "效能分析與平衡測試結果")
hrule(sl, 1.35)

# Chart: bullets per second per difficulty
add_textbox(sl, "各難度每秒產生子彈數量", 0.4, 1.55, 5.6, 0.32,
            font_size=13, bold=True, color=WHITE)

cd = ChartData()
cd.categories = ["Easy", "Normal", "Hard", "Lunatic"]
cd.add_series("普通敵人", [1.2, 1.7, 2.2, 2.8])
cd.add_series("Boss 戰",  [2.5, 3.8, 5.2, 7.0])

chart = sl.shapes.add_chart(
    XL_CHART_TYPE.COLUMN_CLUSTERED,
    Inches(0.4), Inches(1.92), Inches(5.5), Inches(2.6),
    cd
)
chart.chart.plots[0].series[0].format.fill.solid()
chart.chart.plots[0].series[0].format.fill.fore_color.rgb = BLUE
chart.chart.plots[0].series[1].format.fill.solid()
chart.chart.plots[0].series[1].format.fill.fore_color.rgb = RED

# Right: key metrics
metrics = [
    ("最大同時子彈數",   "~ 300 顆",         GOLD),
    ("碰撞偵測耗時",     "< 1ms / frame",    GREEN),
    ("EXP 公式驗證",    "線性成長 ✓",        TEAL),
    ("追蹤彈收斂性",    "3~5 秒內命中 ✓",   ACCTL),
]
for i, (metric, value, clr) in enumerate(metrics):
    y = 1.55 + i * 1.02
    card(sl, 6.1, y, 3.5, 0.87)
    sh = sl.shapes.add_shape(1, Inches(6.1), Inches(y), Inches(0.08), Inches(0.87))
    sh.fill.solid(); sh.fill.fore_color.rgb = clr; sh.line.fill.background()
    add_textbox(sl, metric, 6.28, y + 0.06, 3.2, 0.3, font_size=11, color=DIM)
    add_textbox(sl, value,  6.28, y + 0.38, 3.2, 0.38, font_size=16, bold=True, color=clr)

# Wave offset chart
add_textbox(sl, "波動彈 sin 軌跡（side offset = sin(t)×12）",
            0.4, 4.57, 9.2, 0.25, font_size=10, color=DIM)

import math
cd2 = ChartData()
cd2.categories = [f"t={i*10}" for i in range(10)]
cd2.add_series("side offset (px)", [round(math.sin(i*0.6)*12) for i in range(10)])

chart2 = sl.shapes.add_chart(
    XL_CHART_TYPE.LINE,
    Inches(0.4), Inches(4.82), Inches(9.2), Inches(0.68),
    cd2
)
chart2.chart.plots[0].series[0].format.line.color.rgb = TEAL

# ═══════════════════════════════════════════════════════
# Slide 10 — 未來改進
# ═══════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
slide_bg(sl)
section_tag(sl, "05 未來改進")
slide_title(sl, "演算法優化與擴展方向")
hrule(sl, 1.35)

plans = [
    ("短期目標", GREEN, [
        "Grid Partition 空間分割\n碰撞偵測降至 O(n log n)",
        "Object Pool 物件池\n減少 Bullet GC 壓力",
        "Bezier 曲線彈道\n平滑參數化路徑",
    ]),
    ("中期目標", GOLD, [
        "機率式難度自適應\n依玩家存活率動態調整",
        "Boss AI 狀態機\nIdle/Attack/Enrage 三態",
        "擴充 10 種技能樹\n更深的成長系統",
    ]),
    ("長期目標", RED, [
        "強化學習難度調整\nReinforcement Learning",
        "程序生成 Boss 技能\n隨機但平衡的組合",
        "多人連線彈幕同步\n演算法一致性保障",
    ]),
]
for i, (phase, clr, items) in enumerate(plans):
    x = 0.4 + i * 3.15
    card(sl, x, 1.55, 3.0, 3.75)
    card(sl, x + 0.7, 1.65, 1.6, 0.4, clr)
    add_textbox(sl, phase, x + 0.7, 1.65, 1.6, 0.4,
                font_size=12, bold=True, color=DARK, align=PP_ALIGN.CENTER)
    for j, item in enumerate(items):
        add_oval(sl, x + 0.2, 2.2 + j*0.94, 0.17, 0.17, clr)
        add_textbox(sl, item, x + 0.48, 2.12 + j*0.94, 2.4, 0.82,
                    font_size=11, color=WHITE)

# ═══════════════════════════════════════════════════════
# Slide 11 — 結語
# ═══════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
slide_bg(sl)

add_oval(sl, -1.0, 3.2, 5.5, 5.5, RGBColor(0x3D,0x00,0x66))
add_oval(sl, 7.5, -1.2, 5.5, 5.5, RGBColor(0x00,0x3D,0x66))

add_textbox(sl, "THANK YOU", 0.5, 1.0, 9.0, 1.0,
            font_size=46, bold=True, color=ACCTL, align=PP_ALIGN.CENTER)
add_textbox(sl, "演算法設計篇 完", 0.5, 2.1, 9.0, 0.6,
            font_size=20, color=WHITE, align=PP_ALIGN.CENTER)

sum_stats = [
    ("12",     "彈幕演算法",  "純靜態 BulletPattern 類"),
    ("4",      "子彈運動型",  "直線/追蹤/波動/加速"),
    ("5",      "技能升級樹",  "射速/射頻/生命/炸彈/刷掠"),
    ("O(n²)", "碰撞偵測",    "60fps 下 < 1ms 耗時"),
]
for i, (v, l, sub) in enumerate(sum_stats):
    x = 0.5 + i * 2.28
    card(sl, x, 3.2, 2.1, 1.85)
    add_textbox(sl, v, x, 3.28, 2.1, 0.7,
                font_size=28, bold=True, color=TEAL, align=PP_ALIGN.CENTER)
    add_textbox(sl, l, x, 3.96, 2.1, 0.32,
                font_size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_textbox(sl, sub, x, 4.26, 2.1, 0.62,
                font_size=9, color=DIM, align=PP_ALIGN.CENTER)

# ─── Save ─────────────────────────────────────────────
out = r"C:\project T\TouhouGame\ppt\02_演算法設計.pptx"
prs.save(out)
print(f"✓ Saved: {out}")
