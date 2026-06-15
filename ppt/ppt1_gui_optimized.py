"""
PPT 1 — Bullet Storm：圖形介面設計篇（優化版）
- 精準的 UI 尺寸與布局參數（800×600, 480×560 遊戲區）
- 完整的角色與難度系統介紹
- HUD 面板詳細設計（分數、炸彈、心量、等級、EXP）
- 實際色彩方案與品牌指南
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import math

# ─── Palette (Navy/Cyan theme) ─────────────────────────
DARK   = RGBColor(0x0A, 0x0A, 0x18)
PANEL  = RGBColor(0x10, 0x15, 0x2B)
ACCENT = RGBColor(0x00, 0xC8, 0xD7)
ACCTL  = RGBColor(0x40, 0xE0, 0xFF)
TEAL   = RGBColor(0x00, 0xF5, 0xD4)
GOLD   = RGBColor(0xFF, 0xD7, 0x00)
WHITE  = RGBColor(0xEA, 0xF0, 0xF8)
DIM    = RGBColor(0x80, 0x98, 0xB0)
RED    = RGBColor(0xFF, 0x4D, 0x6D)
GREEN  = RGBColor(0x4D, 0xFF, 0x8C)
BLUE   = RGBColor(0x00, 0xBB, 0xFF)
PANEL2 = RGBColor(0x08, 0x0D, 0x1F)

prs = Presentation()
prs.slide_width  = Inches(10)
prs.slide_height = Inches(5.625)

# ─── Helpers ────────────────────────────────────────────
def slide_bg(slide):
    bg = slide.background; f = bg.fill
    f.solid(); f.fore_color.rgb = DARK

def card(slide, x, y, w, h, fill=None):
    c = fill or PANEL
    sh = slide.shapes.add_shape(5, Inches(x), Inches(y), Inches(w), Inches(h))
    sh.fill.solid(); sh.fill.fore_color.rgb = c
    sh.line.color.rgb = RGBColor(0x20,0x3D,0x5A); sh.line.width = Pt(0.5)
    return sh

def add_textbox(slide, text, x, y, w, h,
                font_size=14, bold=False, color=WHITE,
                align=PP_ALIGN.LEFT, italic=False, font_name="Calibri"):
    txb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = txb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = align
    run = p.add_run()
    run.text = text; run.font.size = Pt(font_size)
    run.font.bold = bold; run.font.italic = italic
    run.font.color.rgb = color; run.font.name = font_name
    return txb

def section_tag(slide, text):
    card(slide, 0.4, 0.15, 2.3, 0.4, ACCENT)
    add_textbox(slide, text, 0.4, 0.15, 2.3, 0.4,
                font_size=11, bold=True, color=DARK, align=PP_ALIGN.CENTER)

def slide_title(slide, text):
    add_textbox(slide, text, 0.4, 0.6, 9.2, 0.7, font_size=28, bold=True, color=WHITE)

def hrule(slide, y):
    sh = slide.shapes.add_shape(1, Inches(0.4), Inches(y), Inches(9.2), Inches(0.03))
    sh.fill.solid(); sh.fill.fore_color.rgb = ACCENT; sh.line.fill.background()

def add_rect_outline(slide, x, y, w, h, color, width=1):
    sh = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(0.02))
    sh.fill.solid(); sh.fill.fore_color.rgb = color
    sh.line.fill.background()
    for (sx, sy, sw, sh_h) in [(x, y, 0.02, h), (x+w-0.02, y, 0.02, h)]:
        s = slide.shapes.add_shape(1, Inches(sx), Inches(sy), Inches(sw), Inches(sh_h))
        s.fill.solid(); s.fill.fore_color.rgb = color; s.line.fill.background()
    sh = slide.shapes.add_shape(1, Inches(x), Inches(y+h-0.02), Inches(w), Inches(0.02))
    sh.fill.solid(); sh.fill.fore_color.rgb = color; sh.line.fill.background()

# ═══════════════════════════════════════════════════════
# Slide 1 — Title
# ═══════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
slide_bg(sl)

# Animated grid background effect
for i in range(12):
    sh = sl.shapes.add_shape(1, Inches(i*0.8-0.5), Inches(2.8), Inches(0.02), Inches(0.3))
    sh.fill.solid(); sh.fill.fore_color.rgb = RGBColor(0, int(100+i*10), 100)
    sh.line.fill.background()

card(sl, 0.5, 0.45, 2.8, 0.52, ACCENT)
add_textbox(sl, "BULLET STORM", 0.5, 0.45, 2.8, 0.52,
            font_size=13, bold=True, color=DARK, align=PP_ALIGN.CENTER)

add_textbox(sl, "圖形介面設計", 0.5, 1.15, 7, 1.2, font_size=46, bold=True, color=WHITE)
add_textbox(sl, "Graphical User Interface Design & Visual System",
            0.5, 2.45, 7, 0.55, font_size=18, color=ACCTL)

sh = sl.shapes.add_shape(1, Inches(0.5), Inches(3.1), Inches(3.5), Inches(0.04))
sh.fill.solid(); sh.fill.fore_color.rgb = ACCENT; sh.line.fill.background()

add_textbox(sl, "800×600 解析度  ·  480×560 遊戲區域  ·  右側 HUD 面板\n2 位角色選擇  ·  4 級難度系統  ·  實時數值反饋",
            0.5, 3.22, 7, 0.85, font_size=14, color=DIM)

stats = [("800×600","視窗尺寸"), ("480×560","遊戲區"), ("2","可選角色"), ("4","難度級別")]
for i, (v, l) in enumerate(stats):
    x = 0.5 + i * 2.3
    card(sl, x, 4.35, 2.1, 0.95)
    add_textbox(sl, v, x, 4.38, 2.1, 0.48, font_size=20, bold=True, color=ACCTL, align=PP_ALIGN.CENTER)
    add_textbox(sl, l, x, 4.84, 2.1, 0.3,  font_size=10, color=DIM,  align=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════
# Slide 2 — Layout 布局
# ═══════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
slide_bg(sl)
section_tag(sl, "01 介面配置")
slide_title(sl, "視窗佈局與座標系統")
hrule(sl, 1.35)

# Draw the layout diagram
layout_x, layout_y = 1.2, 1.55

# Outer frame 800×600
add_rect_outline(sl, layout_x, layout_y, 5.0, 3.75, ACCTL, 2)

# Game field 480×560
field_x_norm = (32/800) * 5.0
field_y_norm = (16/600) * 3.75
field_w_norm = (480/800) * 5.0
field_h_norm = (560/600) * 3.75
add_rect_outline(sl, layout_x + field_x_norm, layout_y + field_y_norm, field_w_norm, field_h_norm, GREEN, 1.5)

# HUD panel
hud_x_norm = (544/800) * 5.0
add_rect_outline(sl, layout_x + hud_x_norm, layout_y + field_y_norm, 5.0 - hud_x_norm - 0.03, field_h_norm, GOLD, 1.5)

# Labels
add_textbox(sl, "800×600 總視窗", layout_x - 0.15, layout_y - 0.35, 1.5, 0.3,
            font_size=10, color=ACCTL, bold=True)
add_textbox(sl, "遊戲區 (480×560)", layout_x + field_x_norm - 0.1, layout_y + 1.2, 2.0, 0.3,
            font_size=10, color=GREEN, bold=True)
add_textbox(sl, "HUD 面板", layout_x + hud_x_norm + 0.05, layout_y + 1.2, 1.5, 0.3,
            font_size=10, color=GOLD, bold=True)

# Spec list
specs = [
    ("視窗", "800×600 固定大小，居中縮放"),
    ("遊戲區", "32,16 起點，480×560 尺寸，雙倍緩衝渲染"),
    ("HUD 面板", "544 起點，包含分數、炸彈、心量、等級"),
    ("背景", "150 顆動畫星，深藍色漸層"),
    ("字體", "MS Gothic / SansSerif，反鋸齒開啟"),
]
for i, (label, desc) in enumerate(specs):
    y = 5.45 + i * 0.32
    add_textbox(sl, label + "：", 0.4, y, 1.1, 0.28, font_size=10, bold=True, color=ACCTL)
    add_textbox(sl, desc, 1.55, y, 8.05, 0.28, font_size=10, color=DIM)

# ═══════════════════════════════════════════════════════
# Slide 3 — 遊戲狀態機
# ═══════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
slide_bg(sl)
section_tag(sl, "02 UI 狀態流")
slide_title(sl, "遊戲狀態與轉移")
hrule(sl, 1.35)

states = [
    ("START", "開始畫面", "標題動畫 + 星場背景", ACCENT, 0.8, 1.6),
    ("CHAR SELECT", "角色選擇", "2 位角色卡片對比", BLUE, 3.3, 1.6),
    ("DIFF SELECT", "難度選擇", "4 級難度色卡", RED, 5.8, 1.6),
    ("PLAYING", "遊戲進行", "即時 HUD 更新", GREEN, 0.8, 3.2),
    ("PAUSED", "暫停", "半透明遮蔽層 + 菜單", GOLD, 3.3, 3.2),
    ("GAME OVER", "遊戲結束", "統計面板 / 重試選項", RED, 5.8, 3.2),
]
for (name, ch, desc, clr, x, y) in states:
    card(sl, x, y, 2.0, 1.32)
    sh = sl.shapes.add_shape(9, Inches(x+0.12), Inches(y+0.1), Inches(1.76), Inches(0.5))
    sh.fill.solid(); sh.fill.fore_color.rgb = clr; sh.line.fill.background()
    add_textbox(sl, ch, x+0.12, y+0.1, 1.76, 0.5,
                font_size=12, bold=True, color=DARK, align=PP_ALIGN.CENTER)
    add_textbox(sl, name, x+0.15, y+0.65, 1.7, 0.22,
                font_size=9, color=DIM, align=PP_ALIGN.CENTER, font_name="Courier New")
    add_textbox(sl, desc, x+0.15, y+0.88, 1.7, 0.4,
                font_size=8.5, color=WHITE, align=PP_ALIGN.CENTER)

add_textbox(sl, "Each state has dedicated rendering & input handling. Smooth transitions with fade/slide animations.",
            0.4, 4.65, 9.2, 0.28, font_size=10, color=DIM)

# ═══════════════════════════════════════════════════════
# Slide 4 — 角色選擇
# ═══════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
slide_bg(sl)
section_tag(sl, "03 角色系統")
slide_title(sl, "2 位可選角色 + 屬性對比")
hrule(sl, 1.35)

chars = [
    ("Aurora Striker", "光彩射手", "快速精準型", "快速移動，精確射擊\n單發威力中等，頻率快", ACCTL, 1.0),
    ("Chronos Engineer", "時間工程師", "爆炸威力型", "移動速度中等，爆炸攻擊\n單發威力強，冷卻較長", GOLD, 5.2),
]
for (name, ch, style, desc, clr, x) in chars:
    card(sl, x, 1.6, 3.7, 3.4)
    card(sl, x + 0.15, 1.75, 3.4, 0.52, clr)
    add_textbox(sl, ch, x + 0.15, 1.75, 3.4, 0.52,
                font_size=16, bold=True, color=DARK, align=PP_ALIGN.CENTER)
    add_textbox(sl, name, x + 0.2, 2.35, 3.3, 0.3,
                font_size=12, color=WHITE, align=PP_ALIGN.CENTER)
    add_textbox(sl, style, x + 0.2, 2.68, 3.3, 0.25,
                font_size=10, color=ACCTL, align=PP_ALIGN.CENTER)
    add_textbox(sl, desc, x + 0.3, 3.05, 3.1, 1.2,
                font_size=10, color=WHITE, align=PP_ALIGN.CENTER)
    # Stats
    stats_data = [
        ("移動", "92%", "75%"),
        ("射速", "快", "中"),
        ("威力", "中", "高"),
        ("特性", "精準", "爆炸"),
    ] if "Aurora" in name else [
        ("移動", "75%", "92%"),
        ("射速", "中", "快"),
        ("威力", "高", "中"),
        ("特性", "爆炸", "精準"),
    ]
    for i, (stat, v1, v2) in enumerate(stats_data):
        y = 4.35 + i * 0.32
        bar_color = clr if (i==0 or i==1 or i==3) else BLUE if i==2 else GOLD
        add_textbox(sl, stat, x + 0.2, y, 0.8, 0.28, font_size=9, color=DIM)
        add_textbox(sl, v1 if x < 3 else v2, x + 1.0, y, 2.4, 0.28, font_size=10, bold=True, color=bar_color, align=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════
# Slide 5 — 難度系統
# ═══════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
slide_bg(sl)
section_tag(sl, "04 難度選擇")
slide_title(sl, "4 級難度與參數調整")
hrule(sl, 1.35)

difficulties = [
    ("Easy", "簡單", "新手友善", "子彈速度 0.65× / 數量少 / 敵人 AI 緩慢", RGBColor(100,200,100), 0.4),
    ("Normal", "普通", "平衡體驗", "子彈速度 0.65× / 數量中等 / 標準難度", RGBColor(100,150,255), 2.65),
    ("Hard", "困難", "進階玩家", "子彈速度 0.78× / 數量多 / 敵人AI激進", RGBColor(255,150,50), 4.9),
    ("Lunatic", "狂亂", "終極挑戰", "子彈速度 1.04× / 數量爆表 / Boss 5 階段", RGBColor(255,50,50), 7.15),
]
for (name, ch, subtitle, desc, clr, x) in difficulties:
    card(sl, x, 1.6, 2.15, 3.4)
    card(sl, x + 0.1, 1.75, 1.95, 0.48, clr)
    add_textbox(sl, ch, x + 0.1, 1.75, 1.95, 0.48,
                font_size=14, bold=True, color=DARK, align=PP_ALIGN.CENTER)
    add_textbox(sl, name, x + 0.15, 2.3, 1.85, 0.28,
                font_size=11, color=WHITE, align=PP_ALIGN.CENTER)
    add_textbox(sl, subtitle, x + 0.15, 2.6, 1.85, 0.22,
                font_size=8.5, color=clr, align=PP_ALIGN.CENTER)
    add_textbox(sl, desc, x + 0.15, 2.9, 1.85, 1.0,
                font_size=8.5, color=DIM, align=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════
# Slide 6 — HUD 面板設計
# ═══════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
slide_bg(sl)
section_tag(sl, "05 HUD 面板")
slide_title(sl, "右側資訊面板詳設計")
hrule(sl, 1.35)

hud_items = [
    ("分數", "SCORE", "即時更新，大字黃金色", GOLD, 1.0),
    ("炸彈", "BOMBS", "圖示 + 數字，可視庫存", RED, 1.95),
    ("心量", "HP / MAX HP", "條形進度條，綠色", GREEN, 2.9),
    ("等級", "LEVEL", "玩家成長指標，藍色", BLUE, 3.85),
    ("EXP", "EXP BAR", "升級進度條，紫色", ACCTL, 4.8),
    ("技能", "5 SKILLS", "星級等級顯示 (0-5)", GOLD, 5.75),
]
for (label, en, desc, clr, x) in hud_items:
    card(sl, x, 1.6, 0.88, 1.3)
    sh = sl.shapes.add_shape(9, Inches(x+0.08), Inches(1.78), Inches(0.72), Inches(0.35))
    sh.fill.solid(); sh.fill.fore_color.rgb = clr; sh.line.fill.background()
    add_textbox(sl, label, x + 0.08, 1.78, 0.72, 0.35,
                font_size=10, bold=True, color=DARK, align=PP_ALIGN.CENTER)
    add_textbox(sl, en, x + 0.08, 2.18, 0.72, 0.22,
                font_size=7.5, color=DIM, align=PP_ALIGN.CENTER, font_name="Courier New")
    add_textbox(sl, desc, x + 0.04, 2.42, 0.8, 0.4,
                font_size=7, color=WHITE, align=PP_ALIGN.CENTER)

# Code boxes
add_textbox(sl, "渲染順序：背景 → 遊戲物件 → 彈幕 → HUD → 除錯資訊",
            0.4, 3.05, 9.2, 0.3, font_size=10, color=ACCTL, bold=True)

hud_code = [
    ("分數顯示", "drawString(\"SCORE: \" + score, x, y);"),
    ("血量條", "fillRect(hpCurrent/maxHp * barWidth, ...);"),
    ("EXP 進度", "fillRect(exp/expNext * barWidth, COLOR_PURPLE);"),
    ("技能星級", "for (5) drawStar(exp.getSkillLevel(i));"),
]
for i, (name, code) in enumerate(hud_code):
    y = 3.5 + i * 0.42
    add_textbox(sl, name + "：", 0.4, y, 1.4, 0.32, font_size=9, bold=True, color=ACCTL)
    card(sl, 1.85, y - 0.02, 7.65, 0.36, PANEL2)
    add_textbox(sl, code, 1.95, y, 7.45, 0.32, font_size=8.5, color=TEAL, font_name="Courier New")

# ═══════════════════════════════════════════════════════
# Slide 7 — 色彩方案
# ═══════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
slide_bg(sl)
section_tag(sl, "06 設計系統")
slide_title(sl, "色彩方案與視覺風格")
hrule(sl, 1.35)

# Color palette
colors_data = [
    ("深藍背景", DARK, "#0A0A18"),
    ("面板灰", PANEL, "#10152B"),
    ("主色 Cyan", ACCENT, "#00C8D7"),
    ("亮 Cyan", ACCTL, "#40E0FF"),
    ("黃金分數", GOLD, "#FFD700"),
    ("生命綠", GREEN, "#4DFF8C"),
    ("危險紅", RED, "#FF4D6D"),
    ("資訊藍", BLUE, "#00BBFF"),
]
for i, (name, clr, hex_code) in enumerate(colors_data):
    col, row = i % 2, i // 2
    x, y = 0.4 + col * 4.85, 1.6 + row * 0.9
    sh = sl.shapes.add_shape(5, Inches(x), Inches(y), Inches(0.5), Inches(0.7))
    sh.fill.solid(); sh.fill.fore_color.rgb = clr; sh.line.fill.background()
    add_textbox(sl, name, x + 0.65, y + 0.08, 3.85, 0.28, font_size=11, bold=True, color=WHITE)
    add_textbox(sl, hex_code, x + 0.65, y + 0.38, 3.85, 0.25,
                font_size=9.5, color=DIM, font_name="Courier New")

# Design principles
add_textbox(sl, "設計原則", 0.4, 4.15, 9.2, 0.28, font_size=12, bold=True, color=ACCTL)
principles = [
    "深色主題：護眼，電競美感，星場背景強化沉浸感",
    "色彩信息學：紅=危險, 綠=收益, 藍=資訊, 黃金=重要分值",
    "對比度優化：所有文字都有 AA 級可及性 WCAG 對比度",
    "動畫提示：HUD 更新時微動畫反饋（閃爍、變色、縮放）",
]
for i, p in enumerate(principles):
    y = 4.45 + i * 0.32
    sh = sl.shapes.add_shape(9, Inches(0.5), Inches(y + 0.06), Inches(0.15), Inches(0.15))
    sh.fill.solid(); sh.fill.fore_color.rgb = ACCTL; sh.line.fill.background()
    add_textbox(sl, p, 0.75, y, 8.85, 0.28, font_size=9.5, color=DIM)

# ═══════════════════════════════════════════════════════
# Slide 8 — 動畫與轉場
# ═══════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
slide_bg(sl)
section_tag(sl, "07 動畫系統")
slide_title(sl, "UI 轉場與動畫效果")
hrule(sl, 1.35)

anims = [
    ("標題動畫", "BULLET STORM 文字旋轉 & 縮放 loop", ACCENT, 0.45, 1.6),
    ("選擇高亮", "菜單項目 hover 時縮放 1.1× & 變色", BLUE, 2.75, 1.6),
    ("血量掉血", "紅色閃爍 + 條形倒計時縮短", RED, 5.05, 1.6),
    ("EXP 升級", "紫色脈衝 + level 數字跳動", ACCTL, 7.35, 1.6),
    ("狀態轉移", "淡出舊畫面 → 淡入新畫面（0.3 秒）", GOLD, 0.45, 2.95),
    ("得分彈出", "浮動文字 + 拋物線上升 + 淡出", GREEN, 2.75, 2.95),
    ("炸彈使用", "圓形擴散波紋 + 視覺震撼特效", RED, 5.05, 2.95),
    ("敵人擊殺", "粒子爆發 + 分數彈出 + 音效", GOLD, 7.35, 2.95),
]
for (name, desc, clr, x, y) in anims:
    card(sl, x, y, 2.15, 1.2)
    sh = sl.shapes.add_shape(9, Inches(x + 0.1), Inches(y + 0.08), Inches(0.6), Inches(0.4))
    sh.fill.solid(); sh.fill.fore_color.rgb = clr; sh.line.fill.background()
    add_textbox(sl, name, x + 0.1, y + 0.08, 0.6, 0.4,
                font_size=8.5, bold=True, color=DARK, align=PP_ALIGN.CENTER)
    add_textbox(sl, desc, x + 0.12, y + 0.52, 1.91, 0.6,
                font_size=7.5, color=DIM, align=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════
# Slide 9 — 無障礙與最佳實踐
# ═══════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
slide_bg(sl)
section_tag(sl, "08 使用者體驗")
slide_title(sl, "無障礙設計與最佳實踐")
hrule(sl, 1.35)

practices = [
    ("對比度", "文字對背景 > 4.5:1 (WCAG AA)", "所有白色文字黑背景檢驗通過"),
    ("字體大小", "主標題 36pt, 正文 14-16pt", "行高 1.5× 提升易讀性"),
    ("色彩不依賴", "紅+圖示、綠+圖示等組合", "色盲玩家仍可理解資訊"),
    ("音效與視覺", "關鍵事件同步音效與動畫", "遊戲進度多感官反饋"),
    ("鍵盤控制", "↑↓←→ 移動, Z 射擊, X 炸彈", "無鼠標即可遊玩"),
    ("遊戲速度", "可調幀率 (60/30 FPS 切換)", "適應不同硬體配置"),
]
for i, (category, spec, note) in enumerate(practices):
    y = 1.6 + i * 0.65
    card(sl, 0.4, y, 9.2, 0.58)
    add_textbox(sl, category, 0.55, y + 0.04, 1.8, 0.22,
                font_size=11, bold=True, color=ACCENT)
    add_textbox(sl, spec, 2.45, y + 0.04, 3.5, 0.22,
                font_size=10, color=WHITE)
    add_textbox(sl, note, 6.0, y + 0.04, 3.55, 0.22,
                font_size=9, color=DIM, italic=True)

# ═══════════════════════════════════════════════════════
# Slide 10 — 結語
# ═══════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
slide_bg(sl)

# Background effect
for i in range(8):
    r = 1.0 + i * 0.6
    sh = sl.shapes.add_shape(9, Inches(5 - r/2), Inches(2.8 - r/2), Inches(r), Inches(r))
    sh.fill.background()
    sh.line.color.rgb = RGBColor(0, int(100 + i*15), int(150 - i*10))
    sh.line.width = Pt(0.8)

add_textbox(sl, "BULLET STORM", 0.5, 0.7, 9.0, 0.8,
            font_size=40, bold=True, color=ACCTL, align=PP_ALIGN.CENTER)
add_textbox(sl, "圖形介面設計篇", 0.5, 1.5, 9.0, 0.5,
            font_size=16, color=ACCENT, align=PP_ALIGN.CENTER)

sum_stats = [
    ("800×600","最佳畫面"),
    ("480×560","遊戲區"),
    ("8 種","HUD 元素"),
    ("4 級","難度調整"),
]
for i, (v, l) in enumerate(sum_stats):
    x = 0.5 + i * 2.28
    card(sl, x, 2.9, 2.1, 1.8)
    add_textbox(sl, v, x, 2.98, 2.1, 0.65,
                font_size=24, bold=True, color=ACCTL, align=PP_ALIGN.CENTER)
    add_textbox(sl, l, x, 3.65, 2.1, 0.28,
                font_size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_textbox(sl, "完美適配", x, 3.95, 2.1, 0.5,
                font_size=8, color=DIM, align=PP_ALIGN.CENTER)

# ─── Save ─────────────────────────────────────────────
out = r"C:\project T\TouhouGame\ppt\01_圖形介面設計.pptx"
prs.save(out)
print("Optimized PPT1 saved")
