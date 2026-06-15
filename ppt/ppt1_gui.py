"""
PPT 1 — Bullet Storm：圖形介面設計篇
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
from pptx.enum.dml import MSO_THEME_COLOR
import pptx.oxml.ns as nsmap
from lxml import etree
import math

# ─── Palette ────────────────────────────────────────────
DARK   = RGBColor(0x0D, 0x1B, 0x2A)
PANEL  = RGBColor(0x1B, 0x2A, 0x3B)
ACCENT = RGBColor(0x00, 0xB4, 0xD8)
ACCD   = RGBColor(0x00, 0x77, 0xA8)
GOLD   = RGBColor(0xF4, 0xA2, 0x61)
WHITE  = RGBColor(0xF0, 0xF4, 0xF8)
DIM    = RGBColor(0x8A, 0xA3, 0xB8)
RED    = RGBColor(0xE6, 0x39, 0x46)
GREEN  = RGBColor(0x57, 0xCC, 0x99)
PURPLE = RGBColor(0x9B, 0x5D, 0xE5)
PANEL2 = RGBColor(0x16, 0x20, 0x32)

W = Inches(10)
H = Inches(5.625)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H

# ─── Helpers ────────────────────────────────────────────
def add_rect(slide, x, y, w, h, fill_rgb, radius_emu=0):
    shape = slide.shapes.add_shape(
        pptx.enum.shapes.MSO_SHAPE_TYPE.RECTANGLE if radius_emu == 0 else 2,  # 2=rounded
        Inches(x), Inches(y), Inches(w), Inches(h)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_rgb
    shape.line.fill.background()
    return shape

def add_textbox(slide, text, x, y, w, h,
                font_size=14, bold=False, color=WHITE,
                align=PP_ALIGN.LEFT, wrap=True, italic=False, font_name="Calibri"):
    txBox = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = font_name
    return txBox

def add_oval(slide, x, y, w, h, fill_rgb):
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    shape = slide.shapes.add_shape(9, Inches(x), Inches(y), Inches(w), Inches(h))  # 9=oval
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_rgb
    shape.line.fill.background()
    return shape

def card(slide, x, y, w, h, fill=None):
    c = fill or PANEL
    sh = slide.shapes.add_shape(5, Inches(x), Inches(y), Inches(w), Inches(h))  # 5=rounded
    sh.fill.solid()
    sh.fill.fore_color.rgb = c
    sh.line.color.rgb = RGBColor(0x30,0x45,0x60)
    sh.line.width = Pt(0.5)
    return sh

def slide_bg(slide):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = DARK

def section_tag(slide, text):
    card(slide, 0.4, 0.15, 2.2, 0.4, ACCENT)
    add_textbox(slide, text, 0.4, 0.15, 2.2, 0.4,
                font_size=11, bold=True, color=DARK, align=PP_ALIGN.CENTER)

def slide_title(slide, text):
    add_textbox(slide, text, 0.4, 0.6, 9.2, 0.7,
                font_size=28, bold=True, color=WHITE)

def hrule(slide, y):
    sh = slide.shapes.add_shape(1, Inches(0.4), Inches(y), Inches(9.2), Inches(0.03))
    sh.fill.solid(); sh.fill.fore_color.rgb = ACCENT
    sh.line.fill.background()

def bullet_dot(slide, x, y, color=ACCENT):
    add_oval(slide, x, y, 0.12, 0.12, color)

# ═══════════════════════════════════════════════════════
# Slide 1 — Title
# ═══════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])  # blank
slide_bg(sl)

# Decorative circles
add_oval(sl, 6.5, -1.2, 6.0, 6.0, ACCD)
add_oval(sl, 7.5, 0.0, 3.5, 3.5, ACCENT)

# Badge
card(sl, 0.5, 0.45, 2.8, 0.52, ACCENT)
add_textbox(sl, "BULLET STORM", 0.5, 0.45, 2.8, 0.52,
            font_size=13, bold=True, color=DARK, align=PP_ALIGN.CENTER)

add_textbox(sl, "圖形介面設計", 0.5, 1.15, 7, 1.2,
            font_size=46, bold=True, color=WHITE)
add_textbox(sl, "Graphical User Interface Design", 0.5, 2.45, 7, 0.55,
            font_size=18, color=ACCENT)

# Divider line
sh = sl.shapes.add_shape(1, Inches(0.5), Inches(3.1), Inches(3.5), Inches(0.04))
sh.fill.solid(); sh.fill.fore_color.rgb = GOLD; sh.line.fill.background()

add_textbox(sl, "彈幕遊戲  ·  Java Swing  ·  純圖形渲染\nA Danmaku Bullet-Hell Shooter",
            0.5, 3.22, 7, 0.85, font_size=14, color=DIM)

# Stats
stats = [("3", "遊戲關卡"), ("12", "彈幕演算法"), ("60+", "穩定 FPS"), ("10″", "視窗寬度")]
for i, (v, l) in enumerate(stats):
    x = 0.5 + i * 2.3
    card(sl, x, 4.35, 2.1, 0.95)
    add_textbox(sl, v, x, 4.38, 2.1, 0.48,
                font_size=24, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
    add_textbox(sl, l, x, 4.84, 2.1, 0.3,
                font_size=10, color=DIM, align=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════
# Slide 2 — 目錄
# ═══════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
slide_bg(sl)
section_tag(sl, "CONTENTS")
slide_title(sl, "簡報目錄")
hrule(sl, 1.35)

toc = [
    ("01", "專案動機",  "為何製作彈幕遊戲？靈感來源與目標"),
    ("02", "系統架構",  "視窗佈局、畫面分區、渲染管線"),
    ("03", "使用技術",  "Java Swing / Graphics2D / 雙緩衝渲染"),
    ("04", "實驗結果",  "HUD 展示、特效、角色繪製成果"),
    ("05", "未來改進",  "精靈圖、動畫系統、UI 美化方向"),
]
for i, (n, t, d) in enumerate(toc):
    y = 1.52 + i * 0.8
    card(sl, 0.4, y, 9.2, 0.65)
    add_oval(sl, 0.55, y + 0.1, 0.45, 0.45, ACCENT)
    add_textbox(sl, n, 0.55, y + 0.1, 0.45, 0.45,
                font_size=12, bold=True, color=DARK, align=PP_ALIGN.CENTER)
    add_textbox(sl, t, 1.12, y + 0.07, 2.5, 0.3,
                font_size=14, bold=True, color=WHITE)
    add_textbox(sl, d, 1.12, y + 0.35, 7.9, 0.25,
                font_size=11, color=DIM)

# ═══════════════════════════════════════════════════════
# Slide 3 — 專案動機
# ═══════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
slide_bg(sl)
section_tag(sl, "01 專案動機")
slide_title(sl, "為什麼製作彈幕遊戲？")
hrule(sl, 1.35)

motives = [
    ("🎮", "遊戲類型挑戰性高",  "彈幕遊戲需精密碰撞偵測與即時渲染，是絕佳技術訓練場景"),
    ("🖥", "純 Java 圖形實作",  "不依賴外部引擎，完全使用 Java Swing/Graphics2D 實現視覺效果"),
    ("⚙", "演算法整合展示",    "12 種彈幕模式結合 RPG 升級系統，展示程式廣度與深度"),
    ("🎓", "學術課程要求",      "涵蓋演算法設計、物件導向設計、GUI 程式設計等多項學習目標"),
]
for i, (icon, title, desc) in enumerate(motives):
    y = 1.55 + i * 0.93
    card(sl, 0.4, y, 5.5, 0.78)
    add_textbox(sl, icon, 0.52, y + 0.14, 0.55, 0.5, font_size=22, align=PP_ALIGN.CENTER)
    add_textbox(sl, title, 1.15, y + 0.07, 4.6, 0.3, font_size=13, bold=True, color=WHITE)
    add_textbox(sl, desc,  1.15, y + 0.38, 4.6, 0.32, font_size=10, color=DIM)

card(sl, 6.1, 1.55, 3.5, 3.7, PANEL2)
add_textbox(sl, "核心目標", 6.2, 1.62, 3.3, 0.38,
            font_size=14, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
goals = ["流暢 60 FPS 遊戲體驗", "直觀的 HUD 資訊顯示",
         "豐富的視覺粒子特效", "角色差異化外觀設計",
         "清晰的難度進度回饋", "EXP 升級視覺化系統"]
for i, g in enumerate(goals):
    add_oval(sl, 6.25, 2.12 + i*0.5, 0.17, 0.17, GOLD)
    add_textbox(sl, g, 6.52, 2.06 + i*0.5, 2.85, 0.3, font_size=11, color=WHITE)

# ═══════════════════════════════════════════════════════
# Slide 4 — 系統架構：視窗佈局
# ═══════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
slide_bg(sl)
section_tag(sl, "02 系統架構")
slide_title(sl, "視窗佈局設計")
hrule(sl, 1.35)

card(sl, 0.4, 1.5, 9.2, 3.85)

# Game field
sh = sl.shapes.add_shape(1, Inches(0.7), Inches(1.72), Inches(5.8), Inches(3.35))
sh.fill.solid(); sh.fill.fore_color.rgb = RGBColor(0x08,0x08,0x20)
sh.line.color.rgb = ACCENT; sh.line.width = Pt(1.0)

# Grid lines on field
for i in range(1, 5):
    ln = sl.shapes.add_shape(1, Inches(0.7), Inches(1.72 + i*0.67), Inches(5.8), Inches(0.01))
    ln.fill.solid(); ln.fill.fore_color.rgb = RGBColor(0x1E,0x2A,0x40)
    ln.line.fill.background()

add_textbox(sl, "遊戲場地\nFIELD  480 × 560 px",
            0.7, 2.8, 5.8, 1.0, font_size=13, color=DIM, align=PP_ALIGN.CENTER)

# HUD panel
sh2 = sl.shapes.add_shape(1, Inches(6.7), Inches(1.72), Inches(2.7), Inches(3.35))
sh2.fill.solid(); sh2.fill.fore_color.rgb = RGBColor(0x0A,0x0A,0x23)
sh2.line.color.rgb = ACCENT; sh2.line.width = Pt(1.0)

hud_rows = [
    ("HIGH SCORE", "9999999"),  ("SCORE", "00142680"),
    ("HEALTH ♥♥♥", ""),         ("BOMBS ★★★", ""),
    ("POWER ■■■■□",""),         ("GRAZE  142", ""),
    ("LEVEL  7", ""),            ("EXP ▓▓▓▓▒▒", ""),
]
for i, (lbl, val) in enumerate(hud_rows):
    add_textbox(sl, lbl, 6.75, 1.82 + i*0.31, 2.0, 0.28, font_size=7.5, color=ACCENT)
    if val:
        add_textbox(sl, val, 8.3, 1.82 + i*0.31, 1.0, 0.28,
                    font_size=7.5, color=WHITE, align=PP_ALIGN.RIGHT)

# Annotations
add_textbox(sl, "● 遊戲場地 (480×560)：敵人/子彈/玩家互動",
            0.4, 5.12, 4.5, 0.3, font_size=10, color=DIM)
add_textbox(sl, "● HUD 面板 (200×560)：分數/生命/炸彈/EXP",
            5.0, 5.12, 4.5, 0.3, font_size=10, color=DIM)

# ═══════════════════════════════════════════════════════
# Slide 5 — 渲染管線
# ═══════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
slide_bg(sl)
section_tag(sl, "02 系統架構")
slide_title(sl, "渲染管線與遊戲迴圈")
hrule(sl, 1.35)

pipeline = [
    ("1", "遊戲邏輯\n更新",    "玩家/敵人/子彈\n狀態更新、碰撞\n偵測、EXP結算", ACCD),
    ("2", "離屏緩衝\n繪製",    "繪製至\nBufferedImage\n避免畫面撕裂",           PURPLE),
    ("3", "場景分層\n渲染",    "背景→粒子→子彈\n→敵人→玩家→UI\n嚴格分層順序",    GOLD),
    ("4", "輸出至\n螢幕",      "緩衝圖像貼至\nJPanel\n震動偏移計算",             GREEN),
]
for i, (n, t, d, c) in enumerate(pipeline):
    x = 0.4 + i * 2.35
    card(sl, x, 1.55, 2.15, 2.85)
    add_oval(sl, x + 0.83, 1.65, 0.5, 0.5, c)
    add_textbox(sl, n, x + 0.83, 1.65, 0.5, 0.5,
                font_size=16, bold=True, color=DARK, align=PP_ALIGN.CENTER)
    add_textbox(sl, t, x + 0.1, 2.25, 1.95, 0.6,
                font_size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_textbox(sl, d, x + 0.1, 2.88, 1.95, 0.88,
                font_size=10, color=DIM, align=PP_ALIGN.CENTER)
    if i < 3:
        arr = sl.shapes.add_shape(1, Inches(x+2.15), Inches(2.9), Inches(0.2), Inches(0.03))
        arr.fill.solid(); arr.fill.fore_color.rgb = ACCENT; arr.line.fill.background()

card(sl, 0.4, 4.56, 9.2, 0.82)
add_textbox(sl, "⏱  javax.swing.Timer @ 16ms 間隔 ≈ 62.5 FPS  |  paintComponent() 每幀更新",
            0.5, 4.6, 9.0, 0.75, font_size=13, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════
# Slide 6 — 使用技術
# ═══════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
slide_bg(sl)
section_tag(sl, "03 使用技術")
slide_title(sl, "圖形渲染技術棧")
hrule(sl, 1.35)

techs = [
    ("核心框架", ACCENT, ["Java Swing — JPanel + JFrame", "javax.swing.Timer 遊戲主迴圈", "ActionListener 每幀回呼"]),
    ("繪圖引擎", GOLD,   ["Graphics2D — 2D 向量繪圖", "RenderingHints 抗鋸齒", "GradientPaint 漸層色彩", "BasicStroke 自訂線條"]),
    ("效能最佳化", GREEN, ["BufferedImage 雙緩衝", "Shape.setClip() 場地裁切", "ArrayList 快速迭代", "removeIf() 批次清除"]),
    ("視覺效果", PURPLE,  ["Particle System 爆炸粒子", "Screen Shake 畫面震動", "Alpha 透明度光暈", "Polygon 角色手繪外型"]),
]
for i, (cat, clr, items) in enumerate(techs):
    col, row = i % 2, i // 2
    x, y = 0.4 + col * 4.7, 1.55 + row * 2.05
    card(sl, x, y, 4.5, 1.85)
    card(sl, x + 0.12, y + 0.1, 1.7, 0.32, clr)
    add_textbox(sl, cat, x + 0.12, y + 0.1, 1.7, 0.32,
                font_size=10, bold=True, color=DARK, align=PP_ALIGN.CENTER)
    for j, item in enumerate(items):
        add_oval(sl, x + 0.18, y + 0.54 + j*0.33, 0.1, 0.1, clr)
        add_textbox(sl, item, x + 0.36, y + 0.47 + j*0.33, 4.0, 0.3, font_size=11, color=WHITE)

# ═══════════════════════════════════════════════════════
# Slide 7 — 實驗結果：HUD 功能
# ═══════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
slide_bg(sl)
section_tag(sl, "04 實驗結果")
slide_title(sl, "HUD 介面設計成果")
hrule(sl, 1.35)

features = [
    ("📊", "分數系統",  "HIGH SCORE + 即時 SCORE，等寬字型確保對齊"),
    ("❤",  "生命顯示",  "圖示化心形顯示，最多 8 格，直觀反映狀態"),
    ("💣", "炸彈顯示",  "★ 星形圖示，升級後炸彈傷害加成視覺標示"),
    ("⚡", "能量條",    "漸層色能量條，4 個力量段位視覺切割"),
    ("🎯", "刷掠計數",  "子彈掠過時累計 Graze 點數，刺激風險操作"),
    ("🌟", "EXP 面板",  "等級、EXP 進度條、5 技能星等即時顯示"),
]
for i, (icon, title, desc) in enumerate(features):
    col, row = i % 2, i // 2
    x, y = 0.4 + col * 4.7, 1.55 + row * 1.3
    card(sl, x, y, 4.5, 1.12)
    add_textbox(sl, icon, x + 0.12, y + 0.28, 0.55, 0.55, font_size=20, align=PP_ALIGN.CENTER)
    add_textbox(sl, title, x + 0.76, y + 0.1, 3.6, 0.35, font_size=13, bold=True, color=ACCENT)
    add_textbox(sl, desc,  x + 0.76, y + 0.46, 3.6, 0.58, font_size=10, color=DIM)

# ═══════════════════════════════════════════════════════
# Slide 8 — 實驗結果：視覺特效
# ═══════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
slide_bg(sl)
section_tag(sl, "04 實驗結果")
slide_title(sl, "視覺效果成果展示")
hrule(sl, 1.35)

effects = [
    ("子彈光暈效果",  "三層同心圓\n外層半透明光暈\n主色圓體+白核心", ACCENT),
    ("粒子爆炸系統",  "擊毀產生 20 粒子\n重力下落\nAlpha 淡出消失", GOLD),
    ("角色手繪外型",  "手工多邊形繪製\nFocus 顯示旋轉\n菱形判定圈",   PURPLE),
]
for i, (title, desc, clr) in enumerate(effects):
    x = 0.4 + i * 3.15
    card(sl, x, 1.55, 3.05, 2.4)
    add_textbox(sl, title, x + 0.1, 1.62, 2.85, 0.38,
                font_size=13, bold=True, color=clr, align=PP_ALIGN.CENTER)
    add_oval(sl, x + 0.9, 2.08, 1.2, 1.2, PANEL2)
    add_textbox(sl, "✦", x + 0.9, 2.08, 1.2, 1.2,
                font_size=28, color=clr, align=PP_ALIGN.CENTER)
    add_textbox(sl, desc, x + 0.1, 3.3, 2.85, 0.58,
                font_size=10, color=DIM, align=PP_ALIGN.CENTER)

# FPS info card
card(sl, 0.4, 4.1, 9.2, 1.2)
add_textbox(sl, "效能表現", 0.55, 4.15, 2, 0.35, font_size=13, bold=True, color=ACCENT)
perf_items = [
    ("最大子彈數量", "~300 顆/幀"),
    ("渲染耗時",    "< 10ms/幀"),
    ("碰撞偵測",   "< 1ms/幀"),
    ("GC 暫停",    "< 5ms"),
]
for i, (label, val) in enumerate(perf_items):
    x = 0.55 + i * 2.25
    add_textbox(sl, label, x, 4.55, 2.0, 0.28, font_size=10, color=DIM)
    add_textbox(sl, val,   x, 4.82, 2.0, 0.38, font_size=18, bold=True, color=WHITE)

# ═══════════════════════════════════════════════════════
# Slide 9 — 未來改進
# ═══════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
slide_bg(sl)
section_tag(sl, "05 未來改進")
slide_title(sl, "圖形介面改進方向")
hrule(sl, 1.35)

plans = [
    ("短期目標", GREEN, [
        "載入外部精靈圖 (Sprite Sheet)\n取代程式繪製角色",
        "背景捲動圖層\n(Parallax Scrolling)",
        "音效與 BGM 整合\n(javax.sound)",
    ]),
    ("中期目標", GOLD, [
        "GPU 加速渲染\n(JavaFX Canvas / JOGL)",
        "自訂字體 TTF 載入\n提升文字美感",
        "動畫狀態機\n角色移動/受傷動畫幀",
    ]),
    ("長期目標", RED, [
        "切換至 LibGDX 框架\n提升跨平台效能",
        "著色器效果 (Shader)\n光暈、閃光、扭曲",
        "支援 4K 解析度\n與多螢幕縮放",
    ]),
]
for i, (phase, clr, items) in enumerate(plans):
    x = 0.4 + i * 3.15
    card(sl, x, 1.55, 3.0, 3.75)
    card(sl, x + 0.7, 1.65, 1.6, 0.4, clr)
    add_textbox(sl, phase, x + 0.7, 1.65, 1.6, 0.4,
                font_size=12, bold=True, color=DARK, align=PP_ALIGN.CENTER)
    for j, item in enumerate(items):
        add_oval(sl, x + 0.2, 2.2 + j*0.94, 0.17, 0.17, clr)
        add_textbox(sl, item, x + 0.48, 2.12 + j*0.94, 2.4, 0.82, font_size=11, color=WHITE)

# ═══════════════════════════════════════════════════════
# Slide 10 — 結語
# ═══════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
slide_bg(sl)

add_oval(sl, -1.0, 3.5, 5.5, 5.5, ACCD)
add_oval(sl, 7.5, -1.0, 4.5, 4.5, ACCENT)

add_textbox(sl, "THANK YOU", 0.5, 1.0, 9.0, 1.0,
            font_size=46, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
add_textbox(sl, "圖形介面設計篇 完", 0.5, 2.1, 9.0, 0.6,
            font_size=20, color=WHITE, align=PP_ALIGN.CENTER)

sum_stats = [
    ("3",   "畫面狀態機",  "START/PLAYING/GAME OVER"),
    ("6",   "HUD 模組",   "分數/生命/炸彈/能量/EXP/Boss"),
    ("∞",   "粒子特效",   "爆炸＋炸彈＋光暈"),
    ("60",  "FPS 目標",   "javax.swing.Timer @16ms"),
]
for i, (v, l, sub) in enumerate(sum_stats):
    x = 0.5 + i * 2.28
    card(sl, x, 3.2, 2.1, 1.85)
    add_textbox(sl, v, x, 3.28, 2.1, 0.7,
                font_size=32, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
    add_textbox(sl, l, x, 3.96, 2.1, 0.32,
                font_size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_textbox(sl, sub, x, 4.26, 2.1, 0.62,
                font_size=9, color=DIM, align=PP_ALIGN.CENTER)

# ─── Save ─────────────────────────────────────────────
out = r"C:\project T\TouhouGame\ppt\01_圖形介面設計.pptx"
prs.save(out)
print(f"✓ Saved: {out}")
